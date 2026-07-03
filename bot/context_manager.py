# c:/Users/adri1/Documents/GitHub/LLM-Discord-Bot/bot/context_manager.py
import discord
import logging
from .store import Store
from .llm_provider import LiteLLMProvider
from .config import Config
import asyncio
import base64
import copy
import importlib.util
import json
import mimetypes
import os
import re
import shutil
import socket
# subprocess is restricted to validated ffmpeg/ffprobe invocations.
import subprocess  # nosec B404
import tempfile
import aiohttp
import PyPDF2
import docx
import openpyxl
from PIL import Image, ImageOps
import io
from ipaddress import ip_address
from datetime import datetime, timezone
from urllib.parse import urlparse

# Optional imports for media processing
try:
    from faster_whisper import WhisperModel
    FASTER_WHISPER_AVAILABLE = True
except ImportError:
    WhisperModel = None
    FASTER_WHISPER_AVAILABLE = False
    logging.warning("faster-whisper not available. Local audio transcription will use fallback engines if available.")

openai_whisper = None
OPENAI_WHISPER_AVAILABLE = importlib.util.find_spec("whisper") is not None

# Optional imports for advanced document processing
try:
    import pytesseract
    try:
        pytesseract.get_tesseract_version()
        TESSERACT_AVAILABLE = True
    except Exception as e:
        TESSERACT_AVAILABLE = False
        logging.warning("pytesseract is installed, but the Tesseract executable is unavailable. OCR features disabled: %s", e)
except ImportError:
    TESSERACT_AVAILABLE = False
    logging.warning("pytesseract not available. OCR features will be limited.")

try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. PowerPoint processing will be limited.")

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    logging.warning("pdfplumber not available. Advanced PDF processing will be limited.")

AUDIO_EXTENSIONS = {
    "aac", "aiff", "alac", "amr", "flac", "m4a", "mp3", "oga", "ogg",
    "opus", "wav", "weba", "wma",
}
VIDEO_EXTENSIONS = {
    "avi", "flv", "m4v", "mkv", "mov", "mp4", "mpeg", "mpg", "ogv",
    "webm", "wmv",
}
TEXT_EXTENSIONS = {
    "txt", "md", "json", "csv", "py", "js", "html", "css", "xml",
    "yaml", "yml", "toml", "ini", "log", "sql",
}
OFFICE_EXTENSIONS = {"docx", "xlsx", "pptx"}

class ContextManager:
    """
    Manages conversation context for the Discord bot.
    
    This class handles:
    - Building context from message history and reply chains
    - Channel summaries and automatic updates
    - User profile management and AI-generated summaries
    - Media processing integration
    - Settings inheritance (guild -> channel overrides)
    """
    def __init__(self, store: Store, llm_provider: LiteLLMProvider, bot=None):
        self.store = store
        self.llm_provider = llm_provider
        self.config = Config()
        self.bot = bot
        self._stt_lock = asyncio.Lock()
        self._faster_whisper_model = None
        self._faster_whisper_model_key = None
        self._openai_whisper_model = None
        self._openai_whisper_model_name = None
        self._guild_emoji_cache = {}
        self._guild_emoji_cache_ttl_seconds = 300

    async def get_guild_and_channel_settings(self, guild_id, channel_id):
        guild_settings = await self.store.get_guild_settings(guild_id)
        channel_overrides = guild_settings.get("channel_overrides", {}).get(str(channel_id), {})
        logging.debug(
            "Loaded guild/channel settings. guild_id=%s channel_id=%s guild_keys=%s override_keys=%s",
            guild_id,
            channel_id,
            sorted(guild_settings.keys()),
            sorted(channel_overrides.keys()) if isinstance(channel_overrides, dict) else [],
        )
        
        final_settings = guild_settings.copy()
        final_settings.update(channel_overrides)

        media_settings = self._default_media_settings()
        media_settings = self._deep_merge_dict(media_settings, guild_settings.get("media", {}))
        media_settings = self._deep_merge_dict(media_settings, channel_overrides.get("media", {}))
        final_settings["media"] = media_settings
        
        return final_settings

    async def is_channel_llm_blacklisted(self, guild_id, channel_id) -> bool:
        settings = await self.get_guild_and_channel_settings(guild_id, channel_id)
        return self.is_llm_blacklisted_settings(settings, channel_id)

    def is_llm_blacklisted_settings(self, settings: dict, channel_id) -> bool:
        if not isinstance(settings, dict):
            return False

        if self._safe_bool(settings.get("llm_blacklisted"), default=False):
            return True

        channel_id = str(channel_id)
        blacklisted_channels = settings.get("llm_blacklisted_channels", [])
        if isinstance(blacklisted_channels, (str, int)):
            blacklisted_channels = [blacklisted_channels]

        return channel_id in {str(value) for value in blacklisted_channels or []}

    def _default_media_settings(self) -> dict:
        return {
            "images": {
                "enabled": self.config.DEFAULT_MEDIA_IMAGES_ENABLED,
                "max_size_mb": 10,
                "ocr_enabled": True,
                "include_ocr_for_vision_models": True,
                "max_ocr_chars": 4000,
                "description_fallback": True,
                "gif": {
                    "extract_frames": True,
                    "max_frames": self.config.DEFAULT_GIF_MAX_FRAMES,
                    "frame_quality": self.config.DEFAULT_GIF_FRAME_QUALITY,
                },
            },
            "audio": {
                "enabled": self.config.DEFAULT_MEDIA_AUDIO_ENABLED,
                "max_size_mb": 25,
                "max_duration_seconds": 300,
                "transcribe": True,
                "include_timestamps": False,
                "max_transcript_chars": 12000,
                "stt_engine": self.config.LOCAL_STT_ENGINE,
                "stt_model": self.config.LOCAL_STT_MODEL,
                "stt_device": self.config.LOCAL_STT_DEVICE,
                "stt_compute_type": self.config.LOCAL_STT_COMPUTE_TYPE,
                "stt_beam_size": self.config.LOCAL_STT_BEAM_SIZE,
                "stt_vad_filter": self.config.LOCAL_STT_VAD_FILTER,
                "stt_language": self.config.LOCAL_STT_LANGUAGE,
            },
            "video": {
                "enabled": self.config.DEFAULT_MEDIA_VIDEO_ENABLED,
                "max_size_mb": self.config.DEFAULT_VIDEO_MAX_SIZE_MB,
                "max_duration_seconds": 120,
                "extract_audio": True,
                "extract_frames": True,
                "frame_interval_seconds": 10,
                "max_frames": self.config.DEFAULT_VIDEO_MAX_FRAMES,
                "frame_quality": self.config.DEFAULT_VIDEO_FRAME_QUALITY,
                "probe_timeout_seconds": self.config.DEFAULT_MEDIA_PROBE_TIMEOUT_SECONDS,
                "frame_timeout_seconds": self.config.DEFAULT_VIDEO_FRAME_TIMEOUT_SECONDS,
                "ocr_frames_for_text_models": True,
                "max_transcript_chars": 12000,
                "stt_engine": self.config.LOCAL_STT_ENGINE,
                "stt_model": self.config.LOCAL_STT_MODEL,
                "stt_device": self.config.LOCAL_STT_DEVICE,
                "stt_compute_type": self.config.LOCAL_STT_COMPUTE_TYPE,
                "stt_beam_size": self.config.LOCAL_STT_BEAM_SIZE,
                "stt_vad_filter": self.config.LOCAL_STT_VAD_FILTER,
                "stt_language": self.config.LOCAL_STT_LANGUAGE,
            },
            "pdf": {
                "enabled": self.config.DEFAULT_MEDIA_PDF_ENABLED,
                "max_size_mb": 10,
                "preserve_formatting": True,
            },
            "office_documents": {
                "enabled": self.config.DEFAULT_MEDIA_OFFICE_DOCUMENTS_ENABLED,
                "max_size_mb": 10,
                "preserve_structure": True,
            },
            "text_files": {
                "enabled": self.config.DEFAULT_MEDIA_TEXT_FILES_ENABLED,
                "max_size_mb": 5,
                "supported_extensions": sorted(TEXT_EXTENSIONS),
            },
            "other_files": {
                "enabled": self.config.DEFAULT_MEDIA_OTHER_FILES_ENABLED,
                "max_size_mb": 20,
                "include_metadata_only": True,
            },
        }

    def _deep_merge_dict(self, base: dict, override: dict) -> dict:
        result = copy.deepcopy(base) if isinstance(base, dict) else {}
        if not isinstance(override, dict):
            return result

        for key, value in override.items():
            if isinstance(value, dict) and isinstance(result.get(key), dict):
                result[key] = self._deep_merge_dict(result[key], value)
            else:
                result[key] = copy.deepcopy(value)
        return result

    async def build_context(self, message: discord.Message = None, channel: discord.TextChannel = None, model_name: str = None, 
                           prompt: str = None, behavior_override: str = None, capabilities_override: str = None,
                           include_bot_identity: bool = True, include_channel_summary: bool = True, 
                           include_user_profiles: bool = True, include_conversation_history: bool = True,
                           include_reply_chain: bool = True, include_current_message: bool = True,
                           include_server_emojis: bool = True):
        """
        Build context for a specific model. If model_name is not provided, uses MAIN_LLM_MODEL.
        
        Args:
            message: Optional. The message to build context for. If not provided, uses the latest message from the channel.
            channel: Optional. The channel to build context for. Required if message is not provided.
            model_name: Optional. The model to build context for. Defaults to MAIN_LLM_MODEL.
            prompt: Optional. Custom prompt to use instead of the current message content. Useful for cogs.
            behavior_override: Optional. Custom behavior prompt to override the default/settings behavior prompt.
            capabilities_override: Optional. Custom capabilities prompt to override the default capabilities prompt.
            include_bot_identity: Optional. Include bot name and ID in system messages. Default True.
            include_channel_summary: Optional. Include channel summary in context. Default True.
            include_user_profiles: Optional. Include user profiles (manual notes + AI summaries). Default True.
            include_conversation_history: Optional. Include recent message history. Default True.
            include_reply_chain: Optional. Include reply chain context (requires include_conversation_history). Default True.
            include_current_message: Optional. Include the current message in history (ignored if prompt is provided). Default True.
            include_server_emojis: Optional. Include custom server emojis available to the bot. Default True.
        
        According to the spec, both decision and main models should receive identical media processing
        but processed according to each model's individual capabilities.
        """
        # Handle the case where no message is provided
        if message is None:
            if channel is None:
                raise ValueError("Either message or channel must be provided")
            
            # Fetch the latest message from the channel
            try:
                async for latest_msg in channel.history(limit=1):
                    message = latest_msg
                    break
                else:
                    # No messages in channel
                    raise ValueError(f"No messages found in channel {channel.name}")
            except Exception as e:
                raise ValueError(f"Could not fetch latest message from channel: {e}")
        
        # If channel wasn't provided but we have a message, get channel from message
        if channel is None:
            channel = message.channel
            
        guild_id = str(channel.guild.id)
        channel_id = str(channel.id)
        user_id = str(message.author.id)

        settings = await self.get_guild_and_channel_settings(guild_id, channel_id)
        data = await self.store.get_guild_data(guild_id)

        # Use specified model or default to main model
        target_model = model_name or self.config.MAIN_LLM_MODEL
        logging.info(
            "Building LLM context. message_id=%s channel_id=%s guild_id=%s author_id=%s model=%s attachments=%s embeds=%s",
            getattr(message, "id", None),
            channel_id,
            guild_id,
            user_id,
            target_model,
            len(getattr(message, "attachments", []) or []),
            len(getattr(message, "embeds", []) or []),
        )
        logging.debug(
            "Context include flags. bot_identity=%s channel_summary=%s user_profiles=%s history=%s reply_chain=%s current_message=%s prompt_override=%s",
            include_bot_identity,
            include_channel_summary,
            include_user_profiles,
            include_conversation_history,
            include_reply_chain,
            include_current_message,
            bool(prompt),
        )

        # 1. System Prompts with optional overrides
        capabilities_prompt = capabilities_override or self.config.CAPABILITIES_PROMPT
        behavior_prompt = behavior_override or settings.get("behavior_prompt", self.config.BEHAVIOR_PROMPT)
        system_prompt = f"{capabilities_prompt}\n\n{behavior_prompt}"

        messages = [{"role": "system", "content": system_prompt}]

        # Developer override: when a trusted developer/operator authored the current message,
        # inject a high-priority instruction so the bot obeys operational commands/meta questions
        # directly instead of deflecting in-character. The prompt is conditional: the bot stays in
        # character for casual chat and only breaks persona for actual requests. Keyed on Discord
        # user ID, which cannot be spoofed.
        if self._is_developer_message(message):
            messages.append({"role": "system", "content": self.config.DEVELOPER_PROMPT})
            logging.info(
                "Developer override applied. message_id=%s author_id=%s",
                getattr(message, "id", None),
                user_id,
            )

        # Add bot identity info as a system message
        if include_bot_identity and self.bot and self.bot.user:
            bot_identity = f"Your Bot name: {self.bot.user.name}\nYour Bot user ID: {self.bot.user.id}"
            messages.append({"role": "system", "content": bot_identity})

        # 2. Server Emojis
        if include_server_emojis:
            emoji_context = await self._server_emoji_context(getattr(channel, "guild", None))
            if emoji_context:
                messages.append({"role": "system", "content": emoji_context})

        # 3. Channel Summary
        if include_channel_summary:
            mentioned_channel_ids = self._message_channel_mention_ids(message)
            channel_data = data.get("channels", {}).get(channel_id, {})
            current_summary = (channel_data.get("summary") or "").strip()
            if current_summary:
                messages.append({
                    "role": "system",
                    "content": self._format_channel_summary_context(
                        channel_name=getattr(channel, "name", channel_id),
                        channel_id=channel_id,
                        summary=current_summary,
                        channel_data=channel_data,
                        preloaded_explicit=channel_id in mentioned_channel_ids,
                        current_channel=True,
                    ),
                })
            messages.extend(
                self._mentioned_channel_summary_messages(
                    message=message,
                    data=data,
                    current_channel_id=channel_id,
                    settings=settings,
                )
            )

        # 4. User Profiles
        if include_user_profiles:
            user_data = data.get("users", {}).get(user_id, {})
            user_profile_content = []
            if "manual_note" in user_data:
                user_profile_content.append(
                    f"Manual note about {message.author.display_name} (User ID: {user_id}): {user_data['manual_note']}"
                )
            if "ai_summary" in user_data:
                user_profile_content.append(
                    f"AI summary of {message.author.display_name} (User ID: {user_id}): {user_data['ai_summary']}"
                )
            
            if user_profile_content:
                messages.append({"role": "system", "content": "\n".join(user_profile_content)})

        reply_chain = []

        # 5. Conversation History (Reply Chain + Recent Messages)
        if include_conversation_history:
            history_limit = settings.get("context", {}).get("history_messages", 15)
            reply_chain_limit = settings.get("context", {}).get("reply_chain_limit", 5)

            # Fetch reply chain
            if include_reply_chain:
                current_message = message
                for _ in range(reply_chain_limit):
                    if current_message.reference and current_message.reference.message_id:
                        try:
                            ref_message = await message.channel.fetch_message(current_message.reference.message_id)
                            reply_chain.insert(0, ref_message)
                            current_message = ref_message
                        except discord.NotFound:
                            break
                    else:
                        break
            
            # Fetch recent messages
            recent_messages = [msg async for msg in message.channel.history(limit=history_limit)]
            recent_messages.reverse() # Oldest to newest
            logging.debug(
                "Fetched context history. message_id=%s reply_chain_count=%s recent_count=%s history_limit=%s reply_chain_limit=%s",
                getattr(message, "id", None),
                len(reply_chain),
                len(recent_messages),
                history_limit,
                reply_chain_limit,
            )

            # Combine and deduplicate
            all_messages = {msg.id: msg for msg in reply_chain}
            all_messages.update({msg.id: msg for msg in recent_messages})
            
            # Always exclude the current message from history to prevent duplication
            # The current message will be added separately at the end
            if message.id in all_messages:
                del all_messages[message.id]

            sorted_messages = sorted(all_messages.values(), key=lambda m: m.created_at)

            if sorted_messages:
                context_lines = [
                    "Recent Discord conversation context, oldest to newest.",
                    "These lines are background only. Do not answer them one by one; answer only the current message that follows.",
                ]
                context_lines.extend(self._format_message_context_line(msg) for msg in sorted_messages)
                messages.append({"role": "system", "content": "\n".join(context_lines)})

        if (
            include_conversation_history
            and include_reply_chain
            and include_current_message
            and not prompt
            and self._should_include_reply_chain_media_context(message, reply_chain)
        ):
            media_context_messages = await self._reply_chain_media_context_messages(reply_chain, target_model)
            messages.extend(media_context_messages)

        # 6. Current Message or Custom Prompt
        if include_current_message or prompt:
            messages.append({
                "role": "system",
                "content": (
                    "Reply only to the next/current Discord message or task. "
                    "Do not prefix the reply with your bot name, username, role label, or user ID."
                )
            })

        if prompt:
            # Use custom prompt instead of message content
            messages.append({"role": "user", "content": prompt})
        else:
            # Add the current message only if include_current_message is True
            if include_current_message:
                current_message_content = await self._format_message_content(
                    message,
                    target_model,
                    current_message=True
                )
                current_message_stats = self._content_stats_for_logging([{"role": "user", "content": current_message_content}])
                if current_message_stats["image_parts"] > 0:
                    messages.append({
                        "role": "system",
                        "content": (
                            "The next/current Discord message includes successfully attached visual input. "
                            "Use the image or extracted video frames when answering visual questions. "
                            "Do not claim the visual input is missing; if visual details are unreadable, say exactly that."
                        ),
                    })
                messages.append({"role": "user", "content": current_message_content})

        logging.info(
            "LLM context built. message_id=%s model=%s messages=%s stats=%s",
            getattr(message, "id", None),
            target_model,
            len(messages),
            self._content_stats_for_logging(messages),
        )
        return messages, settings

    async def _server_emoji_context(self, guild: discord.Guild) -> str:
        if not guild:
            return ""

        emojis = await self._get_guild_emojis(guild)
        if not emojis:
            return ""

        lines = [
            "Available Server Emojis:",
            (
                "These are custom emojis from this Discord server. In replies, use the message_format exactly. "
                "For reaction decisions, return either message_format or reaction_format. "
                "Only use emojis marked status=available."
            ),
        ]

        for emoji in sorted(emojis, key=lambda item: (str(getattr(item, "name", "")).lower(), int(getattr(item, "id", 0) or 0))):
            name = str(getattr(emoji, "name", "") or "").strip()
            emoji_id = getattr(emoji, "id", None)
            if not name or not emoji_id:
                continue

            animated = bool(getattr(emoji, "animated", False))
            status = "available" if getattr(emoji, "available", True) else "unavailable"
            message_format = f"<{'a' if animated else ''}:{name}:{emoji_id}>"
            reaction_format = f"{name}:{emoji_id}"
            emoji_type = "animated" if animated else "static"
            lines.append(
                f"- name={name} id={emoji_id} type={emoji_type} status={status} "
                f"message_format={message_format} reaction_format={reaction_format}"
            )

        return "\n".join(lines) if len(lines) > 2 else ""

    async def _get_guild_emojis(self, guild: discord.Guild) -> list:
        guild_id = str(getattr(guild, "id", ""))
        if not guild_id:
            return []

        now = asyncio.get_running_loop().time()
        cached = self._guild_emoji_cache.get(guild_id)
        if cached and now - cached.get("fetched_at", 0) <= self._guild_emoji_cache_ttl_seconds:
            return list(cached.get("emojis") or [])

        emojis = []
        if hasattr(guild, "fetch_emojis"):
            try:
                emojis = await guild.fetch_emojis()
            except discord.HTTPException as e:
                logging.warning("Could not fetch server emojis for guild %s; using cache fallback: %s", guild_id, e)
                emojis = getattr(guild, "emojis", []) or []
            except Exception as e:
                logging.warning("Unexpected error fetching server emojis for guild %s; using cache fallback: %s", guild_id, e)
                emojis = getattr(guild, "emojis", []) or []
        else:
            emojis = getattr(guild, "emojis", []) or []

        emojis = list(emojis or [])
        self._guild_emoji_cache[guild_id] = {
            "fetched_at": now,
            "emojis": emojis,
        }
        logging.debug("Loaded server emojis for context. guild_id=%s emoji_count=%s", guild_id, len(emojis))
        return emojis

    def _is_developer_message(self, message: discord.Message) -> bool:
        """Return True when the message author is a configured, trusted developer/operator."""
        if not getattr(self.config, "DEVELOPER_OVERRIDE_ENABLED", False):
            return False
        developer_ids = getattr(self.config, "DEVELOPER_USER_IDS", None)
        if not developer_ids:
            return False
        author = getattr(message, "author", None)
        author_id = getattr(author, "id", None)
        if author_id is None:
            return False
        # Never let a bot (including this one) trigger the developer override.
        if getattr(author, "bot", False):
            return False
        return str(author_id) in developer_ids

    def _message_channel_mention_ids(self, message: discord.Message) -> set:
        return {
            str(getattr(channel, "id", ""))
            for channel in (getattr(message, "channel_mentions", []) or [])
            if getattr(channel, "id", None)
        }

    def _mentioned_channel_summary_messages(
        self,
        *,
        message: discord.Message,
        data: dict,
        current_channel_id: str,
        settings: dict,
    ) -> list:
        """Return stored summaries for channels explicitly mentioned in the current Discord message."""
        mentioned_channels = getattr(message, "channel_mentions", []) or []
        if not mentioned_channels:
            return []

        context_settings = settings.get("context", {}) if isinstance(settings, dict) else {}
        if not isinstance(context_settings, dict):
            context_settings = {}
        max_channels = self._safe_int(context_settings.get("mentioned_channel_summary_limit"), default=5, minimum=1, maximum=10)
        max_chars = self._safe_int(context_settings.get("mentioned_channel_summary_max_chars"), default=2500, minimum=500, maximum=8000)
        summaries = []
        seen_channel_ids = {str(current_channel_id)}
        channel_data_by_id = data.get("channels", {})

        for mentioned_channel in mentioned_channels:
            if len(summaries) >= max_channels:
                break

            mentioned_channel_id = str(getattr(mentioned_channel, "id", ""))
            if not mentioned_channel_id or mentioned_channel_id in seen_channel_ids:
                continue
            seen_channel_ids.add(mentioned_channel_id)

            if not self._can_include_mentioned_channel_summary(mentioned_channel, message.guild, getattr(message, "author", None)):
                continue

            channel_data = channel_data_by_id.get(mentioned_channel_id, {})
            summary = (channel_data.get("summary") or "").strip()
            if not summary:
                continue

            channel_name = getattr(mentioned_channel, "name", mentioned_channel_id)
            summary = self._truncate_context_text(summary, max_chars)
            summaries.append({
                "role": "system",
                "content": self._format_channel_summary_context(
                    channel_name=channel_name,
                    channel_id=mentioned_channel_id,
                    summary=summary,
                    channel_data=channel_data,
                    preloaded_explicit=True,
                    current_channel=False,
                ),
            })

        return summaries

    def _format_channel_summary_context(
        self,
        *,
        channel_name: str,
        channel_id: str,
        summary: str,
        channel_data: dict,
        preloaded_explicit: bool,
        current_channel: bool,
    ) -> str:
        if preloaded_explicit:
            label = "Preloaded Explicit Channel Summary"
        elif current_channel:
            label = "Current Channel Summary"
        else:
            label = "Channel Summary"

        metadata_parts = []
        if channel_data.get("last_summary_time"):
            metadata_parts.append(f"last_summary_time={channel_data['last_summary_time']}")
        if channel_data.get("messages_since_summary") is not None:
            metadata_parts.append(f"messages_since_summary={channel_data['messages_since_summary']}")
        if channel_data.get("summary_type"):
            metadata_parts.append(f"summary_type={channel_data['summary_type']}")

        metadata = f"\nSummary metadata: {', '.join(metadata_parts)}" if metadata_parts else ""
        instruction = ""
        if preloaded_explicit:
            instruction = (
                "\nThis summary was preloaded because the current message explicitly mentions this channel. "
                "Use it directly for broad recaps; do not call get_channel_summary for this channel unless the user asks for a fresh tool check."
            )

        return f"{label} for #{channel_name} (Channel ID: {channel_id}):\n{summary}{metadata}{instruction}"

    def _can_include_mentioned_channel_summary(
        self,
        channel: discord.abc.GuildChannel,
        current_guild: discord.Guild,
        requester: discord.Member = None,
    ) -> bool:
        if not channel or not current_guild:
            return False

        guild = getattr(channel, "guild", None)
        if not guild or guild.id != current_guild.id:
            return False

        me = getattr(guild, "me", None)
        if not me and self.bot and self.bot.user and hasattr(guild, "get_member"):
            me = guild.get_member(self.bot.user.id)
        if not me or not hasattr(channel, "permissions_for"):
            return False

        try:
            permissions = channel.permissions_for(me)
        except Exception:
            return False

        bot_can_read = bool(
            getattr(permissions, "view_channel", False)
            and getattr(permissions, "read_message_history", False)
        )
        if not bot_can_read:
            return False

        if not requester:
            return False

        if not isinstance(requester, discord.Member) or getattr(getattr(requester, "guild", None), "id", None) != guild.id:
            requester = guild.get_member(getattr(requester, "id", 0)) if hasattr(guild, "get_member") else None
        if not requester:
            return False

        try:
            requester_permissions = channel.permissions_for(requester)
        except Exception:
            return False

        return bool(
            getattr(requester_permissions, "view_channel", False)
            and getattr(requester_permissions, "read_message_history", False)
        )

    def _should_include_reply_chain_media_context(self, message: discord.Message, reply_chain: list) -> bool:
        if not reply_chain:
            return False

        if getattr(message, "attachments", None) or getattr(message, "embeds", None):
            return False

        if not any(self._message_has_media(item) for item in reply_chain):
            return False

        content = (getattr(message, "content", "") or "").lower()
        media_terms = {
            "attachment",
            "attached",
            "file",
            "image",
            "picture",
            "photo",
            "screenshot",
            "video",
            "clip",
            "frame",
            "frames",
            "media",
            "visual",
            "pdf",
            "document",
        }
        return any(term in content for term in media_terms)

    def _message_has_media(self, message: discord.Message) -> bool:
        return bool(getattr(message, "attachments", None) or getattr(message, "embeds", None))

    async def _reply_chain_media_context_messages(self, reply_chain: list, target_model: str) -> list:
        media_messages = [item for item in (reply_chain or []) if self._message_has_media(item)]
        media_messages = media_messages[-2:]
        if not media_messages:
            return []

        context_messages = [{
            "role": "system",
            "content": (
                "The current user is referring to visual/media content from the reply chain. "
                "The following media-bearing Discord message(s) are background context, not new requests. "
                "Use their attached images or extracted video frames when answering the current message."
            ),
        }]

        for media_message in media_messages:
            try:
                processed_content = await self._format_message_content(
                    media_message,
                    target_model,
                    current_message=False,
                )
            except Exception as e:
                logging.exception("Could not process reply-chain media context. message_id=%s", getattr(media_message, "id", None))
                context_messages.append({
                    "role": "system",
                    "content": (
                        f"[Could not process referenced media from message {getattr(media_message, 'id', 'unknown')}: "
                        f"{str(e)[:80]}]"
                    ),
                })
                continue

            context_messages.append({"role": "user", "content": processed_content})

        stats = self._content_stats_for_logging(context_messages)
        logging.info(
            "Included reply-chain media context. source_messages=%s stats=%s",
            [getattr(item, "id", None) for item in media_messages],
            stats,
        )
        return context_messages

    def _safe_int(self, value, *, default: int, minimum: int, maximum: int) -> int:
        try:
            parsed = int(value)
        except (TypeError, ValueError):
            parsed = default
        return max(minimum, min(parsed, maximum))

    def _safe_float(self, value, *, default: float, minimum: float, maximum: float) -> float:
        try:
            parsed = float(value)
        except (TypeError, ValueError):
            parsed = default
        return max(minimum, min(parsed, maximum))

    def _safe_bool(self, value, *, default: bool) -> bool:
        if isinstance(value, bool):
            return value
        if value is None:
            return default
        value = str(value).strip().lower()
        if value in {"true", "1", "yes", "y", "on"}:
            return True
        if value in {"false", "0", "no", "n", "off"}:
            return False
        return default

    def _truncate_document_text(self, text: str, section_config: dict) -> str:
        """Cap extracted document text so a large file cannot blow up the LLM context."""
        max_chars = self._safe_int(
            (section_config or {}).get("max_context_chars"),
            default=20000,
            minimum=1000,
            maximum=200000,
        )
        return self._truncate_context_text(text, max_chars)

    def _truncate_context_text(self, text: str, max_chars: int) -> str:
        text = str(text or "").strip()
        if len(text) <= max_chars:
            return text
        return text[: max_chars - 3].rstrip() + "..."

    def _format_message_context_line(self, message: discord.Message) -> str:
        """
        Format a Discord message as plain background context.

        History is intentionally kept out of user/assistant turns so the model
        does not treat older messages as prompts waiting for answers.
        """
        author_name = getattr(message.author, "display_name", getattr(message.author, "name", "Unknown User"))
        author_id = getattr(message.author, "id", "unknown")
        author_marker = " [bot]" if getattr(message.author, "bot", False) else ""
        timestamp = message.created_at.strftime("%Y-%m-%d %H:%M:%S UTC")
        content = (message.content or "[empty message]").replace("\r", " ").replace("\n", " ").strip()

        extra_parts = []
        if message.attachments:
            filenames = ", ".join(att.filename for att in message.attachments)
            extra_parts.append(f"attachments: {filenames}")
        if message.embeds:
            extra_parts.append(f"embeds: {len(message.embeds)}")
        if message.reference and message.reference.message_id:
            extra_parts.append(f"reply_to_message_id: {message.reference.message_id}")

        extra = f" ({'; '.join(extra_parts)})" if extra_parts else ""
        return f"[{timestamp}] {author_name}{author_marker} (User ID: {author_id}): {content}{extra}"

    async def _format_message_content(self, message: discord.Message, model_name: str = None, *, current_message: bool = False):
        """
        Format message content with attachments processed for the specified model.
        If model_name is not provided, uses MAIN_LLM_MODEL.
        """
        logging.info(
            "Formatting Discord message content. message_id=%s model=%s current_message=%s text_chars=%s attachments=%s embeds=%s",
            getattr(message, "id", None),
            model_name or self.config.MAIN_LLM_MODEL,
            current_message,
            len(message.content or ""),
            len(getattr(message, "attachments", []) or []),
            len(getattr(message, "embeds", []) or []),
        )
        content_parts = []
        
        # Add user information for identification
        user_info = f"{message.author.display_name} (User ID: {message.author.id})"
        prefix = "Current Discord message to answer from " if current_message else "Discord message from "
        if message.content:
            content_parts.append({"type": "text", "text": f"{prefix}{user_info}:\n{message.content}"})
        else:
            content_parts.append({"type": "text", "text": f"{prefix}{user_info}:\n[empty message]"})


        if message.attachments:
            # Use specified model or default to main model for attachment processing
            target_model = model_name or self.config.MAIN_LLM_MODEL
            media_settings = (await self.get_guild_and_channel_settings(message.guild.id, message.channel.id)).get("media", {})
            
            for attachment in message.attachments:
                logging.info(
                    "Processing message attachment. message_id=%s filename=%s content_type=%s declared_size=%s model=%s",
                    getattr(message, "id", None),
                    getattr(attachment, "filename", "unknown"),
                    getattr(attachment, "content_type", None),
                    getattr(attachment, "size", None),
                    target_model,
                )
                processed_content = await self._process_attachment(attachment, target_model, media_settings)
                logging.info(
                    "Attachment processing finished. message_id=%s filename=%s summary=%s",
                    getattr(message, "id", None),
                    getattr(attachment, "filename", "unknown"),
                    self._processed_content_summary(processed_content),
                )
                if isinstance(processed_content, dict) and processed_content.get("type") == "animated_gif":
                    # Handle animated GIF: convert to list of content parts for message processing
                    frame_info = f"Animated GIF: {processed_content['filename']}"
                    total_frames = processed_content.get("total_frames", "unknown")
                    extracted_frames = processed_content.get("extracted_frames", 0)
                    if total_frames > extracted_frames:
                        frame_info += f" (showing {extracted_frames} representative frames from {total_frames} total frames)"
                    else:
                        frame_info += f" ({extracted_frames} frames)"
                    
                    content_parts.append({"type": "text", "text": frame_info})
                    ocr_text = (processed_content.get("ocr_text") or "").strip()
                    if ocr_text:
                        content_parts.append({
                            "type": "text",
                            "text": f"--- OCR text extracted from animated GIF frames: {processed_content['filename']} ---\n{ocr_text}",
                        })
                    for frame_data in processed_content["frames"]:
                        content_parts.append({
                            "type": "image_url", 
                            "image_url": {"url": frame_data}
                        })
                elif isinstance(processed_content, list):
                    # This is a list of content parts (legacy or other processing)
                    content_parts.extend(processed_content)
                elif isinstance(processed_content, dict):
                    # This is a single structured content (like image_url)
                    content_parts.append(processed_content)
                else:
                    # This is a text fallback
                    content_parts.append({"type": "text", "text": processed_content})

        # Process embeds (important for Tenor GIFs and other embedded media)
        if message.embeds:
            target_model = model_name or self.config.MAIN_LLM_MODEL
            media_settings = (await self.get_guild_and_channel_settings(message.guild.id, message.channel.id)).get("media", {})
            
            for embed in message.embeds:
                logging.info(
                    "Processing message embed. message_id=%s embed_type=%s title_present=%s url_present=%s",
                    getattr(message, "id", None),
                    getattr(embed, "type", None),
                    bool(getattr(embed, "title", None)),
                    bool(getattr(embed, "url", None)),
                )
                processed_embed = await self._process_embed(embed, target_model, media_settings)
                logging.info(
                    "Embed processing finished. message_id=%s summary=%s",
                    getattr(message, "id", None),
                    self._processed_content_summary(processed_embed),
                )
                if isinstance(processed_embed, list):
                    content_parts.extend(processed_embed)
                elif isinstance(processed_embed, dict):
                    content_parts.append(processed_embed)
                elif processed_embed:  # Only add non-empty text
                    content_parts.append({"type": "text", "text": processed_embed})
        
        # If we only have text content, return it as a string for simplicity
        if len(content_parts) == 1 and content_parts[0].get("type") == "text":
            logging.debug(
                "Formatted message content as plain text. message_id=%s text_chars=%s",
                getattr(message, "id", None),
                len(content_parts[0]["text"]),
            )
            return content_parts[0]["text"]
        elif content_parts:
            logging.info(
                "Formatted message content as structured parts. message_id=%s summary=%s",
                getattr(message, "id", None),
                self._processed_content_summary(content_parts),
            )
            return content_parts
        else:
            return "[empty message]"

    def _processed_content_summary(self, content) -> str:
        if isinstance(content, str):
            return f"text chars={len(content)} preview={content[:120]!r}"
        if isinstance(content, dict):
            content_type = content.get("type")
            if content_type == "animated_gif":
                return (
                    "animated_gif "
                    f"filename={content.get('filename')} total_frames={content.get('total_frames')} "
                    f"extracted_frames={content.get('extracted_frames')}"
                )
            if content_type in {"image_url", "input_image"}:
                image_url = (content.get("image_url") or {}).get("url") or ""
                source = "inline_base64" if isinstance(image_url, str) and image_url.startswith("data:image/") else "url"
                return f"{content_type} source={source}"
            return f"dict type={content_type or 'unknown'} keys={sorted(content.keys())}"
        if isinstance(content, list):
            stats = self._content_stats_for_logging([{"role": "user", "content": content}])
            return f"list parts={len(content)} stats={stats}"
        if content is None:
            return "none"
        return f"{type(content).__name__}"

    def _content_stats_for_logging(self, messages: list) -> dict:
        stats = {
            "text_chars": 0,
            "image_parts": 0,
            "inline_image_parts": 0,
            "other_parts": 0,
            "roles": {},
        }

        for message in messages or []:
            if isinstance(message, dict):
                role = message.get("role", "unknown")
                stats["roles"][role] = stats["roles"].get(role, 0) + 1
                self._add_content_stats_for_logging(message.get("content"), stats)
            else:
                stats["other_parts"] += 1

        return stats

    def _add_content_stats_for_logging(self, content, stats: dict):
        if isinstance(content, str):
            stats["text_chars"] += len(content)
            return
        if isinstance(content, list):
            for part in content:
                self._add_content_stats_for_logging(part, stats)
            return
        if isinstance(content, dict):
            content_type = content.get("type")
            if content_type == "text":
                stats["text_chars"] += len(content.get("text") or "")
            elif content_type in {"image_url", "input_image"}:
                stats["image_parts"] += 1
                image_url = (content.get("image_url") or {}).get("url") or ""
                if isinstance(image_url, str) and image_url.startswith("data:image/"):
                    stats["inline_image_parts"] += 1
            else:
                stats["other_parts"] += 1
            return
        if content is not None:
            stats["other_parts"] += 1

    def _is_audio_file(self, mime_type: str, file_ext: str) -> bool:
        return bool((mime_type and mime_type.startswith("audio/")) or file_ext in AUDIO_EXTENSIONS)

    def _is_video_file(self, mime_type: str, file_ext: str) -> bool:
        return bool((mime_type and mime_type.startswith("video/")) or file_ext in VIDEO_EXTENSIONS)

    def _image_content_part(self, url: str, file_bytes: bytes = None, mime_type: str = None, model_name: str = None, filename: str = None):
        image_url = url
        if self.llm_provider.prefers_inline_image_data(model_name):
            if file_bytes is None:
                return f"[Image omitted because validated image bytes were unavailable: {filename or url or 'unknown image'}]"
            inline_image = self._normalized_inline_image_data_url(
                file_bytes=file_bytes,
                mime_type=mime_type,
                filename=filename,
            )
            if not inline_image:
                return f"[Image omitted because the downloaded image data was invalid: {filename or url or 'unknown image'}]"
            image_url = inline_image

        return {"type": "image_url", "image_url": {"url": image_url}}

    def _normalized_inline_image_data_url(self, *, file_bytes: bytes, mime_type: str = None, filename: str = None) -> str | None:
        """Validate and re-encode inline images before handing them to strict local providers."""
        try:
            normalized_bytes, normalized_mime = self._normalize_image_bytes_for_inline(file_bytes, mime_type)
        except Exception as e:
            logging.warning(
                "Invalid image data omitted before LLM request. filename=%s mime_type=%s bytes=%s error_type=%s error=%s",
                filename,
                mime_type,
                len(file_bytes or b""),
                type(e).__name__,
                e,
            )
            return None

        encoded_image = base64.b64encode(normalized_bytes).decode("ascii")
        return f"data:{normalized_mime};base64,{encoded_image}"

    def _normalize_image_bytes_for_inline(self, file_bytes: bytes, mime_type: str = None) -> tuple[bytes, str]:
        if not file_bytes:
            raise ValueError("empty image payload")

        with Image.open(io.BytesIO(file_bytes)) as image:
            source_format = (image.format or "").upper()
            image.seek(0)
            image.load()
            image = ImageOps.exif_transpose(image)
            if image.width <= 0 or image.height <= 0:
                raise ValueError(f"invalid image dimensions {image.width}x{image.height}")

            compatible_image = self._image_to_compatible_rgb(image)
            output_mime = self._inline_image_output_mime(mime_type, source_format)
            buffer = io.BytesIO()
            if output_mime == "image/png":
                compatible_image.save(buffer, format="PNG")
            else:
                compatible_image.save(buffer, format="JPEG", quality=92, optimize=True)
            return buffer.getvalue(), output_mime

    def _image_to_compatible_rgb(self, image: Image.Image) -> Image.Image:
        if image.mode in {"RGBA", "LA"} or (image.mode == "P" and "transparency" in image.info):
            rgba_image = image.convert("RGBA")
            background = Image.new("RGBA", rgba_image.size, (255, 255, 255, 255))
            background.alpha_composite(rgba_image)
            return background.convert("RGB")

        if image.mode != "RGB":
            return image.convert("RGB")

        return image.copy()

    def _inline_image_output_mime(self, mime_type: str = None, source_format: str = None) -> str:
        source_mime = (mime_type or "").split(";", 1)[0].strip().lower()
        source_format = (source_format or "").upper()
        if source_mime in {"image/png", "image/gif"} or source_format in {"PNG", "GIF"}:
            return "image/png"
        return "image/jpeg"

    def _vision_image_content_parts(
        self,
        *,
        url: str,
        file_bytes: bytes,
        mime_type: str,
        model_name: str,
        filename: str,
        images_config: dict,
    ):
        image_part = self._image_content_part(url, file_bytes, mime_type, model_name, filename)
        if not isinstance(image_part, dict):
            return image_part

        parts = [
            {
                "type": "text",
                "text": (
                    f"Attached image: {filename or 'image'}. "
                    "Use this visual attachment when answering image-related questions."
                ),
            }
        ]

        ocr_text = self._ocr_image_text_for_context(file_bytes, filename, images_config)
        if ocr_text:
            parts.append({
                "type": "text",
                "text": f"--- OCR text extracted from attached image: {filename or 'image'} ---\n{ocr_text}",
            })

        parts.append(image_part)
        return parts

    def _ocr_image_text_for_context(self, file_bytes: bytes, filename: str, images_config: dict) -> str:
        if not (
            images_config.get("ocr_enabled", True)
            and images_config.get("include_ocr_for_vision_models", True)
            and TESSERACT_AVAILABLE
            and file_bytes
        ):
            return ""

        try:
            with Image.open(io.BytesIO(file_bytes)) as img:
                img = ImageOps.exif_transpose(img)
                ocr_text = pytesseract.image_to_string(img).strip()
        except Exception as e:
            logging.info(
                "Image OCR context extraction skipped. filename=%s error_type=%s error=%s",
                filename,
                type(e).__name__,
                e,
            )
            return ""

        if not ocr_text:
            logging.debug("Image OCR context extraction found no text. filename=%s", filename)
            return ""

        max_chars = self._safe_int(images_config.get("max_ocr_chars"), default=4000, minimum=500, maximum=12000)
        ocr_text = self._truncate_context_text(ocr_text, max_chars)
        logging.info("Image OCR context extracted. filename=%s chars=%s", filename, len(ocr_text))
        return ocr_text

    def _ocr_image_frames_text_for_context(self, frames: list, filename: str, images_config: dict) -> str:
        if not (
            images_config.get("ocr_enabled", True)
            and images_config.get("include_ocr_for_vision_models", True)
            and TESSERACT_AVAILABLE
            and frames
        ):
            return ""

        lines = []
        for frame in frames:
            try:
                with Image.open(io.BytesIO(frame.get("jpeg_bytes") or b"")) as img:
                    text = pytesseract.image_to_string(img).strip()
                if text:
                    lines.append(f"[{frame.get('label') or 'frame'}] {text}")
            except Exception as e:
                logging.debug(
                    "GIF frame OCR skipped. filename=%s label=%s error_type=%s error=%s",
                    filename,
                    frame.get("label"),
                    type(e).__name__,
                    e,
                )

        if not lines:
            logging.debug("GIF frame OCR found no text. filename=%s frames=%s", filename, len(frames or []))
            return ""

        max_chars = self._safe_int(images_config.get("max_ocr_chars"), default=4000, minimum=500, maximum=12000)
        ocr_text = self._truncate_context_text("\n".join(lines), max_chars)
        logging.info("GIF frame OCR context extracted. filename=%s frames=%s chars=%s", filename, len(frames), len(ocr_text))
        return ocr_text

    async def _write_temp_file(self, file_bytes: bytes, file_ext: str) -> str:
        suffix = f".{file_ext.lstrip('.')}" if file_ext else ""
        path = await asyncio.to_thread(self._write_temp_file_sync, file_bytes, suffix)
        logging.debug("Wrote temporary media file. path=%s bytes=%s suffix=%s", path, len(file_bytes or b""), suffix)
        return path

    def _write_temp_file_sync(self, file_bytes: bytes, suffix: str) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
            temp_file.write(file_bytes)
            return temp_file.name

    def _remove_temp_file(self, path: str):
        if not path:
            return
        try:
            os.remove(path)
            logging.debug("Removed temporary media file. path=%s", path)
        except FileNotFoundError:
            pass
        except OSError as e:
            logging.warning(f"Could not remove temporary media file {path}: {e}")

    async def _validate_public_media_url(self, url: str) -> str | None:
        parsed = urlparse(url or "")
        if parsed.scheme not in {"http", "https"}:
            return "Only http and https attachment URLs are allowed."
        if not parsed.hostname:
            return "Attachment URL must include a hostname."
        if parsed.username or parsed.password:
            return "Attachment URLs with embedded credentials are not allowed."

        host = parsed.hostname.strip().lower()
        if host == "localhost" or host.endswith(".local"):
            return "Local or private attachment hostnames are not allowed."

        try:
            literal_ip = ip_address(host)
            if self._is_private_or_local_ip(literal_ip):
                return "Local or private attachment IP addresses are not allowed."
            return None
        except ValueError:
            pass

        loop = asyncio.get_running_loop()
        try:
            infos = await loop.getaddrinfo(host, parsed.port or (443 if parsed.scheme == "https" else 80), type=socket.SOCK_STREAM)
        except socket.gaierror:
            return "Could not resolve attachment hostname."

        for info in infos:
            resolved_host = info[4][0]
            try:
                resolved_ip = ip_address(resolved_host)
            except ValueError:
                continue
            if self._is_private_or_local_ip(resolved_ip):
                return "Attachment URL resolves to a local or private network address."

        return None

    def _is_private_or_local_ip(self, value) -> bool:
        return any([
            value.is_private,
            value.is_loopback,
            value.is_link_local,
            value.is_multicast,
            value.is_reserved,
            value.is_unspecified,
        ])

    def _attachment_download_limit_bytes(self, media_settings: dict, mime_type: str, file_ext: str, is_likely_gif: bool) -> int:
        if (mime_type and mime_type.startswith("image/")) or is_likely_gif:
            max_size_mb = self._safe_float(media_settings.get("images", {}).get("max_size_mb"), default=10.0, minimum=0.1, maximum=100.0)
        elif self._is_audio_file(mime_type, file_ext):
            max_size_mb = self._safe_float(media_settings.get("audio", {}).get("max_size_mb"), default=25.0, minimum=0.1, maximum=500.0)
        elif self._is_video_file(mime_type, file_ext):
            max_size_mb = self._safe_float(media_settings.get("video", {}).get("max_size_mb"), default=250.0, minimum=0.1, maximum=1000.0)
        elif mime_type == "application/pdf":
            max_size_mb = self._safe_float(media_settings.get("pdf", {}).get("max_size_mb"), default=10.0, minimum=0.1, maximum=100.0)
        elif file_ext in OFFICE_EXTENSIONS:
            max_size_mb = self._safe_float(media_settings.get("office_documents", {}).get("max_size_mb"), default=10.0, minimum=0.1, maximum=100.0)
        elif (mime_type and mime_type.startswith("text/")) or file_ext in TEXT_EXTENSIONS:
            max_size_mb = self._safe_float(media_settings.get("text_files", {}).get("max_size_mb"), default=5.0, minimum=0.1, maximum=100.0)
        else:
            max_size_mb = self._safe_float(media_settings.get("other_files", {}).get("max_size_mb"), default=20.0, minimum=0.1, maximum=100.0)

        return max(1, int(max_size_mb * 1024 * 1024))

    async def _read_bounded_response_body(self, resp, max_bytes: int) -> tuple[bytes, bool]:
        chunks = []
        total_bytes = 0
        read_limit = max(1, int(max_bytes)) + 1

        async for chunk in resp.content.iter_chunked(64 * 1024):
            if not chunk:
                continue

            remaining = read_limit - total_bytes
            if remaining > 0:
                chunks.append(chunk[:remaining])
            total_bytes += len(chunk)

            if total_bytes > max_bytes:
                return b"".join(chunks), True

        return b"".join(chunks), False

    async def _get_media_duration(self, file_path: str, media_kind: str, media_config: dict = None):
        timeout_seconds = self._safe_float(
            (media_config or {}).get("probe_timeout_seconds"),
            default=self.config.DEFAULT_MEDIA_PROBE_TIMEOUT_SECONDS,
            minimum=1.0,
            maximum=120.0,
        )
        logging.debug("Reading %s duration. file=%s timeout=%.1fs", media_kind, file_path, timeout_seconds)
        return await asyncio.to_thread(self._get_media_duration_sync, file_path, media_kind, timeout_seconds)

    def _get_media_duration_sync(self, file_path: str, media_kind: str, timeout_seconds: float = 10.0):
        duration = self._probe_duration_with_ffprobe(file_path, timeout_seconds)
        if duration is not None:
            logging.debug("Read %s duration with ffprobe. file=%s duration=%s", media_kind, file_path, duration)
            return duration

        duration = self._probe_duration_with_ffmpeg(file_path, timeout_seconds)
        if duration is not None:
            logging.debug("Read %s duration with ffmpeg metadata. file=%s duration=%s", media_kind, file_path, duration)
            return duration

        logging.warning("Could not read %s duration with bounded FFmpeg probes. file=%s", media_kind, file_path)
        return None

    def _probe_duration_with_ffprobe(self, file_path: str, timeout_seconds: float):
        ffprobe = self._ffprobe_executable()
        if not ffprobe:
            return None

        command = [
            ffprobe,
            "-v",
            "error",
            "-show_entries",
            "format=duration",
            "-of",
            "json",
            file_path,
        ]
        try:
            completed = self._run_media_subprocess(command, timeout_seconds)
            payload = json.loads(completed.stdout or "{}")
            duration = payload.get("format", {}).get("duration")
            if duration not in (None, "N/A", ""):
                return max(0.0, float(duration))
        except subprocess.TimeoutExpired:
            logging.warning("ffprobe duration probe timed out after %.1fs. file=%s", timeout_seconds, file_path)
        except Exception as e:
            logging.debug("ffprobe duration probe failed. file=%s error=%s", file_path, e)
        return None

    def _probe_duration_with_ffmpeg(self, file_path: str, timeout_seconds: float):
        ffmpeg = self._ffmpeg_executable()
        if not ffmpeg:
            return None

        command = [
            ffmpeg,
            "-hide_banner",
            "-nostdin",
            "-i",
            file_path,
        ]
        try:
            completed = self._run_media_subprocess(command, timeout_seconds)
            output = f"{completed.stderr or ''}\n{completed.stdout or ''}"
            match = re.search(r"Duration:\s*(\d+):(\d+):(\d+(?:\.\d+)?)", output)
            if not match:
                return None
            hours, minutes, seconds = match.groups()
            return int(hours) * 3600 + int(minutes) * 60 + float(seconds)
        except subprocess.TimeoutExpired:
            logging.warning("ffmpeg duration metadata probe timed out after %.1fs. file=%s", timeout_seconds, file_path)
        except Exception as e:
            logging.debug("ffmpeg duration metadata probe failed. file=%s error=%s", file_path, e)
        return None

    def _run_media_subprocess(self, command: list, timeout_seconds: float):
        command = self._validate_media_subprocess_command(command)
        logging.debug("Running media subprocess. command=%s timeout=%.1fs", command, timeout_seconds)
        startupinfo = None
        creationflags = 0
        if os.name == "nt":
            startupinfo = subprocess.STARTUPINFO()
            startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
            creationflags = getattr(subprocess, "CREATE_NO_WINDOW", 0)

        # Command executable is validated against ffmpeg/ffprobe and shell=False is used.
        return subprocess.run(  # nosec B603
            command,
            capture_output=True,
            text=True,
            timeout=timeout_seconds,
            stdin=subprocess.DEVNULL,
            startupinfo=startupinfo,
            creationflags=creationflags,
        )

    def _run_media_subprocess_bytes(self, command: list, timeout_seconds: float):
        command = self._validate_media_subprocess_command(command)
        logging.debug("Running media bytes subprocess. command=%s timeout=%.1fs", command, timeout_seconds)
        startupinfo = None
        creationflags = 0
        if os.name == "nt":
            startupinfo = subprocess.STARTUPINFO()
            startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
            creationflags = getattr(subprocess, "CREATE_NO_WINDOW", 0)

        # Command executable is validated against ffmpeg/ffprobe and shell=False is used.
        return subprocess.run(  # nosec B603
            command,
            capture_output=True,
            timeout=timeout_seconds,
            stdin=subprocess.DEVNULL,
            startupinfo=startupinfo,
            creationflags=creationflags,
        )

    def _validate_media_subprocess_command(self, command: list) -> list:
        if not isinstance(command, list) or not command:
            raise ValueError("Media subprocess command must be a non-empty list.")
        if not all(isinstance(part, str) for part in command):
            raise ValueError("Media subprocess command parts must be strings.")

        executable = self._normalize_executable_path(command[0])
        allowed_executables = {
            normalized
            for normalized in (
                self._normalize_executable_path(self._ffmpeg_executable()),
                self._normalize_executable_path(self._ffprobe_executable()),
            )
            if normalized
        }
        if executable not in allowed_executables:
            raise ValueError(f"Media subprocess executable is not allowed: {command[0]}")

        return [os.path.abspath(command[0]), *command[1:]]

    def _normalize_executable_path(self, path: str | None) -> str | None:
        if not path:
            return None
        return os.path.normcase(os.path.abspath(path))

    def _ffprobe_executable(self):
        ffprobe = shutil.which("ffprobe")
        if ffprobe:
            return ffprobe

        ffmpeg = shutil.which("ffmpeg")
        if ffmpeg:
            sibling = os.path.join(os.path.dirname(ffmpeg), "ffprobe.exe" if os.name == "nt" else "ffprobe")
            if os.path.exists(sibling):
                return sibling
        return None

    def _ffmpeg_executable(self):
        ffmpeg = shutil.which("ffmpeg")
        if ffmpeg:
            return ffmpeg

        try:
            import imageio_ffmpeg
            return imageio_ffmpeg.get_ffmpeg_exe()
        except Exception as e:
            logging.debug("imageio_ffmpeg fallback unavailable: %s", e)
            return None

    async def _transcribe_media_file(self, file_path: str, media_config: dict) -> dict:
        logging.info(
            "Starting local media transcription. file=%s engine=%s model=%s device=%s",
            file_path,
            media_config.get("stt_engine") or self.config.LOCAL_STT_ENGINE,
            media_config.get("stt_model") or self.config.LOCAL_STT_MODEL,
            media_config.get("stt_device") or self.config.LOCAL_STT_DEVICE,
        )
        async with self._stt_lock:
            transcript = await asyncio.to_thread(self._transcribe_media_file_sync, file_path, media_config)
        logging.info(
            "Finished local media transcription. file=%s engine=%s chars=%s segments=%s language=%s",
            file_path,
            transcript.get("engine"),
            len(transcript.get("text") or ""),
            len(transcript.get("segments") or []),
            transcript.get("language"),
        )
        return transcript

    def _transcribe_media_file_sync(self, file_path: str, media_config: dict) -> dict:
        engine = str(media_config.get("stt_engine") or self.config.LOCAL_STT_ENGINE).strip().lower().replace("_", "-")

        if engine in {"faster", "faster-whisper"}:
            if FASTER_WHISPER_AVAILABLE:
                return self._transcribe_with_faster_whisper(file_path, media_config)
            if OPENAI_WHISPER_AVAILABLE:
                logging.warning("faster-whisper is unavailable; falling back to openai-whisper local transcription.")
                return self._transcribe_with_openai_whisper(file_path, media_config)
            raise RuntimeError("No local speech-to-text engine is installed. Install faster-whisper.")

        if engine in {"openai-whisper", "whisper"}:
            if OPENAI_WHISPER_AVAILABLE:
                return self._transcribe_with_openai_whisper(file_path, media_config)
            if FASTER_WHISPER_AVAILABLE:
                logging.warning("openai-whisper is unavailable; falling back to faster-whisper local transcription.")
                return self._transcribe_with_faster_whisper(file_path, media_config)
            raise RuntimeError("No local speech-to-text engine is installed. Install faster-whisper.")

        raise RuntimeError(f"Unsupported local speech-to-text engine: {engine}")

    def _transcribe_with_faster_whisper(self, file_path: str, media_config: dict) -> dict:
        model_name = str(media_config.get("stt_model") or self.config.LOCAL_STT_MODEL)
        device = str(media_config.get("stt_device") or self.config.LOCAL_STT_DEVICE)
        compute_type = str(media_config.get("stt_compute_type") or self.config.LOCAL_STT_COMPUTE_TYPE)
        beam_size = self._safe_int(
            media_config.get("stt_beam_size", self.config.LOCAL_STT_BEAM_SIZE),
            default=self.config.LOCAL_STT_BEAM_SIZE,
            minimum=1,
            maximum=10,
        )
        vad_filter = self._safe_bool(
            media_config.get("stt_vad_filter", self.config.LOCAL_STT_VAD_FILTER),
            default=self.config.LOCAL_STT_VAD_FILTER,
        )
        language = media_config.get("stt_language", self.config.LOCAL_STT_LANGUAGE)
        language = str(language).strip() if language else None

        model = self._get_faster_whisper_model(model_name, device, compute_type)
        kwargs = {
            "beam_size": beam_size,
            "vad_filter": vad_filter,
        }
        if language:
            kwargs["language"] = language

        segments, info = model.transcribe(file_path, **kwargs)
        segment_list = [
            {"start": segment.start, "end": segment.end, "text": segment.text.strip()}
            for segment in segments
        ]
        text = " ".join(segment["text"] for segment in segment_list if segment["text"]).strip()
        logging.debug(
            "faster-whisper transcription details. file=%s segments=%s chars=%s language=%s probability=%s",
            file_path,
            len(segment_list),
            len(text),
            getattr(info, "language", None),
            getattr(info, "language_probability", None),
        )
        return {
            "engine": "faster-whisper",
            "model": model_name,
            "language": getattr(info, "language", None),
            "language_probability": getattr(info, "language_probability", None),
            "segments": segment_list,
            "text": text,
        }

    def _get_faster_whisper_model(self, model_name: str, device: str, compute_type: str):
        key = (model_name, device, compute_type)
        if self._faster_whisper_model is None or self._faster_whisper_model_key != key:
            logging.info(f"Loading faster-whisper model '{model_name}' on {device} with compute_type={compute_type}.")
            self._faster_whisper_model = WhisperModel(model_name, device=device, compute_type=compute_type)
            self._faster_whisper_model_key = key
        return self._faster_whisper_model

    def _transcribe_with_openai_whisper(self, file_path: str, media_config: dict) -> dict:
        model_name = str(media_config.get("stt_model") or self.config.LOCAL_STT_MODEL)
        beam_size = self._safe_int(
            media_config.get("stt_beam_size", self.config.LOCAL_STT_BEAM_SIZE),
            default=self.config.LOCAL_STT_BEAM_SIZE,
            minimum=1,
            maximum=10,
        )
        language = media_config.get("stt_language", self.config.LOCAL_STT_LANGUAGE)
        language = str(language).strip() if language else None

        model = self._get_openai_whisper_model(model_name)
        kwargs = {"fp16": False, "beam_size": beam_size}
        if language:
            kwargs["language"] = language
        result = model.transcribe(file_path, **kwargs)
        segment_list = [
            {
                "start": segment.get("start", 0),
                "end": segment.get("end", 0),
                "text": str(segment.get("text", "")).strip(),
            }
            for segment in result.get("segments", [])
        ]
        text = str(result.get("text", "")).strip()
        if not text:
            text = " ".join(segment["text"] for segment in segment_list if segment["text"]).strip()
        logging.debug(
            "openai-whisper transcription details. file=%s segments=%s chars=%s language=%s",
            file_path,
            len(segment_list),
            len(text),
            result.get("language"),
        )
        return {
            "engine": "openai-whisper-local",
            "model": model_name,
            "language": result.get("language"),
            "language_probability": None,
            "segments": segment_list,
            "text": text,
        }

    def _get_openai_whisper_model(self, model_name: str):
        global openai_whisper
        if openai_whisper is None:
            import whisper as openai_whisper_module
            openai_whisper = openai_whisper_module

        if self._openai_whisper_model is None or self._openai_whisper_model_name != model_name:
            logging.info(f"Loading openai-whisper fallback model '{model_name}'.")
            self._openai_whisper_model = openai_whisper.load_model(model_name)
            self._openai_whisper_model_name = model_name
        return self._openai_whisper_model

    def _format_transcript_block(self, title: str, filename: str, transcript: dict, duration, media_config: dict) -> str:
        include_timestamps = self._safe_bool(media_config.get("include_timestamps"), default=False)
        max_chars = self._safe_int(
            media_config.get("max_transcript_chars", 12000),
            default=12000,
            minimum=500,
            maximum=50000,
        )

        metadata = []
        if duration is not None:
            metadata.append(f"Duration: {self._format_seconds(duration)}")
        if transcript.get("engine"):
            metadata.append(f"Engine: {transcript['engine']}")
        if transcript.get("model"):
            metadata.append(f"STT Model: {transcript['model']}")
        if transcript.get("language"):
            language = transcript["language"]
            probability = transcript.get("language_probability")
            if isinstance(probability, (int, float)):
                metadata.append(f"Detected Language: {language} ({probability:.2f})")
            else:
                metadata.append(f"Detected Language: {language}")

        if include_timestamps and transcript.get("segments"):
            body = "\n".join(
                f"[{self._format_seconds(segment['start'])} -> {self._format_seconds(segment['end'])}] {segment['text']}"
                for segment in transcript["segments"]
                if segment.get("text")
            )
        else:
            body = transcript.get("text", "").strip()

        if not body:
            body = "[No speech detected.]"

        body = self._truncate_context_text(body, max_chars)
        header = f"--- {title}: {filename} ---"
        if metadata:
            return f"{header}\n" + "\n".join(metadata) + f"\n\n{body}"
        return f"{header}\n{body}"

    def _format_seconds(self, seconds) -> str:
        try:
            seconds = max(0.0, float(seconds))
        except (TypeError, ValueError):
            return "unknown"

        hours = int(seconds // 3600)
        minutes = int((seconds % 3600) // 60)
        secs = seconds % 60
        if hours:
            return f"{hours:02d}:{minutes:02d}:{secs:06.3f}"
        return f"{minutes:02d}:{secs:06.3f}"

    async def _extract_video_frames(self, file_path: str, video_config: dict) -> dict:
        if not self._ffmpeg_executable():
            raise RuntimeError("ffmpeg is not installed or not available on PATH")
        logging.info(
            "Starting video frame extraction. file=%s max_frames=%s interval=%s quality=%s",
            file_path,
            video_config.get("max_frames", self.config.DEFAULT_VIDEO_MAX_FRAMES),
            video_config.get("frame_interval_seconds", 10),
            video_config.get("frame_quality", self.config.DEFAULT_VIDEO_FRAME_QUALITY),
        )
        return await asyncio.to_thread(self._extract_video_frames_sync, file_path, video_config)

    def _extract_video_frames_sync(self, file_path: str, video_config: dict) -> dict:
        ffmpeg = self._ffmpeg_executable()
        duration = float(self._get_media_duration_sync(
            file_path,
            "video",
            self._safe_float(
                video_config.get("probe_timeout_seconds"),
                default=self.config.DEFAULT_MEDIA_PROBE_TIMEOUT_SECONDS,
                minimum=1.0,
                maximum=120.0,
            ),
        ) or 0)
        max_frames = self._safe_int(video_config.get("max_frames"), default=self.config.DEFAULT_VIDEO_MAX_FRAMES, minimum=1, maximum=30)
        frame_interval = self._safe_float(video_config.get("frame_interval_seconds"), default=10.0, minimum=0.1, maximum=3600.0)
        frame_quality = self._safe_int(video_config.get("frame_quality"), default=self.config.DEFAULT_VIDEO_FRAME_QUALITY, minimum=1, maximum=100)
        frame_timeout = self._safe_float(
            video_config.get("frame_timeout_seconds"),
            default=self.config.DEFAULT_VIDEO_FRAME_TIMEOUT_SECONDS,
            minimum=1.0,
            maximum=120.0,
        )
        frame_times = self._sample_video_times(duration, max_frames, frame_interval)
        logging.info(
            "Video frame sampling plan. file=%s duration=%s max_frames=%s interval=%s quality=%s frame_timeout=%.1fs frame_times=%s",
            file_path,
            self._format_seconds(duration),
            max_frames,
            frame_interval,
            frame_quality,
            frame_timeout,
            [self._format_seconds(t) for t in frame_times],
        )

        frames = []
        qscale = self._jpeg_quality_to_ffmpeg_qscale(frame_quality)
        for frame_time in frame_times:
            command = [
                ffmpeg,
                "-hide_banner",
                "-loglevel",
                "error",
                "-nostdin",
                "-ss",
                f"{frame_time:.3f}",
                "-i",
                file_path,
                "-frames:v",
                "1",
                "-an",
                "-f",
                "image2pipe",
                "-vcodec",
                "mjpeg",
                "-q:v",
                str(qscale),
                "pipe:1",
            ]
            try:
                completed = self._run_media_subprocess_bytes(command, frame_timeout)
                frame_bytes = completed.stdout or b""
                if completed.returncode != 0 or not frame_bytes:
                    stderr = (completed.stderr or b"").decode("utf-8", errors="ignore")
                    logging.warning(
                        "Could not extract video frame. file=%s timestamp=%s returncode=%s stderr=%s",
                        file_path,
                        self._format_seconds(frame_time),
                        completed.returncode,
                        stderr[:500],
                    )
                    continue

                frame_b64 = base64.b64encode(frame_bytes).decode()
                frames.append({
                    "timestamp": frame_time,
                    "data_url": f"data:image/jpeg;base64,{frame_b64}",
                    "jpeg_bytes": frame_bytes,
                })
                logging.debug(
                    "Extracted video frame. file=%s timestamp=%s jpeg_bytes=%s",
                    file_path,
                    self._format_seconds(frame_time),
                    len(frame_bytes),
                )
            except subprocess.TimeoutExpired:
                logging.warning(
                    "Video frame extraction timed out. file=%s timestamp=%s timeout=%.1fs",
                    file_path,
                    self._format_seconds(frame_time),
                    frame_timeout,
                )
            except Exception as e:
                logging.warning("Could not extract video frame at %.2fs from %s: %s", frame_time, file_path, e)

        logging.info(
            "Video frame extraction finished. file=%s extracted=%s requested=%s",
            file_path,
            len(frames),
            len(frame_times),
        )
        return {"duration": duration, "frames": frames}

    def _jpeg_quality_to_ffmpeg_qscale(self, frame_quality: int) -> int:
        quality = max(1, min(100, int(frame_quality)))
        return max(2, min(31, round((100 - quality) / 3.5) + 2))

    def _sample_video_times(self, duration: float, max_frames: int, frame_interval: float) -> list:
        if duration <= 0:
            return [0.0]

        safe_end = max(0.0, duration - 0.1)
        interval_times = []
        current = 0.0
        while current <= safe_end:
            interval_times.append(round(current, 3))
            current += frame_interval

        if interval_times and len(interval_times) <= max_frames:
            return interval_times

        if max_frames == 1:
            return [round(min(safe_end, duration / 2), 3)]

        step = safe_end / (max_frames - 1)
        return [round(min(safe_end, i * step), 3) for i in range(max_frames)]

    async def _ocr_video_frames(self, frames: list) -> str:
        if not TESSERACT_AVAILABLE:
            logging.info("Skipping video frame OCR because pytesseract is unavailable. frames=%s", len(frames or []))
            return ""
        logging.info("Starting video frame OCR. frames=%s", len(frames or []))
        return await asyncio.to_thread(self._ocr_video_frames_sync, frames)

    def _ocr_video_frames_sync(self, frames: list) -> str:
        lines = []
        for frame in frames:
            try:
                with Image.open(io.BytesIO(frame["jpeg_bytes"])) as img:
                    text = pytesseract.image_to_string(img).strip()
                if text:
                    lines.append(f"[{self._format_seconds(frame['timestamp'])}] {text}")
                    logging.debug(
                        "Video frame OCR text found. timestamp=%s chars=%s",
                        self._format_seconds(frame["timestamp"]),
                        len(text),
                    )
            except Exception as e:
                logging.warning(f"Could not OCR video frame at {frame.get('timestamp', 0)}s: {e}")
        logging.info("Video frame OCR finished. frames_with_text=%s chars=%s", len(lines), len("\n".join(lines)))
        return "\n".join(lines)

    async def _process_audio_attachment(
        self,
        attachment: discord.Attachment,
        file_bytes: bytes,
        file_ext: str,
        file_size_mb: float,
        model_name: str,
        audio_config: dict,
    ):
        logging.info(
            "Audio attachment processing started. filename=%s size_mb=%.2f ext=%s model=%s config=%s",
            getattr(attachment, "filename", "unknown"),
            file_size_mb,
            file_ext,
            model_name,
            {
                "enabled": audio_config.get("enabled"),
                "max_size_mb": audio_config.get("max_size_mb"),
                "max_duration_seconds": audio_config.get("max_duration_seconds"),
                "transcribe": audio_config.get("transcribe"),
                "stt_engine": audio_config.get("stt_engine"),
                "stt_model": audio_config.get("stt_model"),
            },
        )
        if not self._safe_bool(audio_config.get("enabled"), default=True):
            logging.info("Audio processing disabled. filename=%s", attachment.filename)
            return f"[Audio processing disabled: {attachment.filename}]"

        max_size = self._safe_float(audio_config.get("max_size_mb"), default=25.0, minimum=0.1, maximum=500.0)
        if file_size_mb > max_size:
            logging.warning(
                "Audio attachment over size limit. filename=%s size_mb=%.2f max_size_mb=%.2f",
                attachment.filename,
                file_size_mb,
                max_size,
            )
            return f"[Audio too large for processing: {attachment.filename} - {file_size_mb:.1f}MB]"

        if self.llm_provider.supports_audio(model_name) and self._safe_bool(audio_config.get("send_direct_if_supported"), default=False):
            logging.info("Sending audio directly to model. filename=%s model=%s", attachment.filename, model_name)
            return {"type": "audio_url", "audio_url": {"url": attachment.url}}

        if not self._safe_bool(audio_config.get("transcribe"), default=True):
            logging.info("Audio transcription disabled; returning metadata only. filename=%s", attachment.filename)
            return f"--- Audio Metadata: {attachment.filename} ---\nFile Size: {file_size_mb:.2f} MB\nTranscription disabled."

        temp_path = await self._write_temp_file(file_bytes, file_ext)
        try:
            duration = await self._get_media_duration(temp_path, "audio", audio_config)
            max_duration = self._safe_float(audio_config.get("max_duration_seconds"), default=300.0, minimum=0.0, maximum=86400.0)
            logging.info(
                "Audio attachment duration checked. filename=%s duration=%s max_duration=%s",
                attachment.filename,
                self._format_seconds(duration),
                self._format_seconds(max_duration),
            )
            if duration is not None and max_duration > 0 and duration > max_duration:
                logging.warning(
                    "Audio attachment over duration limit. filename=%s duration=%s max_duration=%s",
                    attachment.filename,
                    self._format_seconds(duration),
                    self._format_seconds(max_duration),
                )
                return (
                    f"[Audio too long for processing: {attachment.filename} - "
                    f"{self._format_seconds(duration)} > {self._format_seconds(max_duration)}]"
                )

            transcript = await self._transcribe_media_file(temp_path, audio_config)
            logging.info(
                "Audio attachment processing finished. filename=%s transcript_chars=%s segments=%s",
                attachment.filename,
                len(transcript.get("text") or ""),
                len(transcript.get("segments") or []),
            )
            return self._format_transcript_block("Audio Transcript", attachment.filename, transcript, duration, audio_config)
        except Exception as e:
            logging.exception("Audio transcription failed. filename=%s", attachment.filename)
            return f"[Audio transcription failed: {attachment.filename} - {str(e)[:80]}]"
        finally:
            self._remove_temp_file(temp_path)

    async def _process_video_attachment(
        self,
        attachment: discord.Attachment,
        file_bytes: bytes,
        file_ext: str,
        file_size_mb: float,
        model_name: str,
        video_config: dict,
    ):
        logging.info(
            "Video attachment processing started. filename=%s size_mb=%.2f ext=%s model=%s supports_vision=%s config=%s",
            getattr(attachment, "filename", "unknown"),
            file_size_mb,
            file_ext,
            model_name,
            self.llm_provider.supports_vision(model_name),
            {
                "enabled": video_config.get("enabled"),
                "max_size_mb": video_config.get("max_size_mb"),
                "max_duration_seconds": video_config.get("max_duration_seconds"),
                "extract_audio": video_config.get("extract_audio"),
                "extract_frames": video_config.get("extract_frames"),
                "frame_interval_seconds": video_config.get("frame_interval_seconds"),
                "max_frames": video_config.get("max_frames"),
                "frame_quality": video_config.get("frame_quality"),
                "ocr_frames_for_text_models": video_config.get("ocr_frames_for_text_models"),
            },
        )
        if not self._safe_bool(video_config.get("enabled"), default=True):
            logging.info("Video processing disabled. filename=%s", attachment.filename)
            return f"[Video processing disabled: {attachment.filename}]"

        max_size = self._safe_float(video_config.get("max_size_mb"), default=250.0, minimum=0.1, maximum=1000.0)
        if file_size_mb > max_size:
            logging.warning(
                "Video attachment over size limit. filename=%s size_mb=%.2f max_size_mb=%.2f",
                attachment.filename,
                file_size_mb,
                max_size,
            )
            return f"[Video too large for processing: {attachment.filename} - {file_size_mb:.1f}MB]"

        temp_path = await self._write_temp_file(file_bytes, file_ext)
        try:
            duration = await self._get_media_duration(temp_path, "video", video_config)
            max_duration = self._safe_float(video_config.get("max_duration_seconds"), default=120.0, minimum=0.0, maximum=86400.0)
            logging.info(
                "Video attachment duration checked. filename=%s duration=%s max_duration=%s",
                attachment.filename,
                self._format_seconds(duration),
                self._format_seconds(max_duration),
            )
            if duration is not None and max_duration > 0 and duration > max_duration:
                logging.warning(
                    "Video attachment over duration limit. filename=%s duration=%s max_duration=%s",
                    attachment.filename,
                    self._format_seconds(duration),
                    self._format_seconds(max_duration),
                )
                return (
                    f"[Video too long for processing: {attachment.filename} - "
                    f"{self._format_seconds(duration)} > {self._format_seconds(max_duration)}]"
                )

            content_parts = []
            extract_audio = self._safe_bool(video_config.get("extract_audio"), default=True)
            extract_frames = self._safe_bool(video_config.get("extract_frames"), default=True)
            logging.info(
                "Video extraction choices. filename=%s extract_audio=%s extract_frames=%s",
                attachment.filename,
                extract_audio,
                extract_frames,
            )

            if extract_audio:
                try:
                    transcript = await self._transcribe_media_file(temp_path, video_config)
                    content_parts.append({
                        "type": "text",
                        "text": self._format_transcript_block("Video Audio Transcript", attachment.filename, transcript, duration, video_config),
                    })
                    logging.info(
                        "Video audio transcript added. filename=%s transcript_chars=%s segments=%s",
                        attachment.filename,
                        len(transcript.get("text") or ""),
                        len(transcript.get("segments") or []),
                    )
                except Exception as e:
                    logging.exception("Video audio transcription failed. filename=%s", attachment.filename)
                    content_parts.append({
                        "type": "text",
                        "text": f"[Video audio transcription failed: {attachment.filename} - {str(e)[:80]}]",
                    })

            if extract_frames:
                try:
                    frame_result = await self._extract_video_frames(temp_path, video_config)
                    frames = frame_result.get("frames", [])
                    timestamps = ", ".join(self._format_seconds(frame["timestamp"]) for frame in frames)
                    frame_summary = (
                        f"--- Video Frame Samples: {attachment.filename} ---\n"
                        f"Duration: {self._format_seconds(duration if duration is not None else frame_result.get('duration'))}\n"
                        f"Extracted {len(frames)} representative frame(s). "
                        "Each following image is an actual frame extracted from the video."
                    )
                    if timestamps:
                        frame_summary += f" at: {timestamps}"

                    if frames and self.llm_provider.supports_vision(model_name):
                        content_parts.append({"type": "text", "text": frame_summary})
                        for frame in frames:
                            content_parts.append({
                                "type": "text",
                                "text": f"Video frame at {self._format_seconds(frame['timestamp'])} from {attachment.filename}:",
                            })
                            content_parts.append({"type": "image_url", "image_url": {"url": frame["data_url"]}})
                        logging.info(
                            "Video frames attached to vision model request. filename=%s frames=%s model=%s",
                            attachment.filename,
                            len(frames),
                            model_name,
                        )
                    elif frames:
                        ocr_text = ""
                        if self._safe_bool(video_config.get("ocr_frames_for_text_models"), default=True):
                            ocr_text = await self._ocr_video_frames(frames)

                        if ocr_text:
                            logging.info(
                                "Video frame OCR added for text-only model. filename=%s frames=%s ocr_chars=%s model=%s",
                                attachment.filename,
                                len(frames),
                                len(ocr_text),
                                model_name,
                            )
                            content_parts.append({
                                "type": "text",
                                "text": f"--- Video Frame OCR: {attachment.filename} ---\n{ocr_text}",
                            })
                        else:
                            logging.info(
                                "Video frames extracted but not attached; target model has no vision and OCR produced no text. filename=%s frames=%s model=%s",
                                attachment.filename,
                                len(frames),
                                model_name,
                            )
                            content_parts.append({
                                "type": "text",
                                "text": frame_summary + "\n[Frames were extracted locally but not attached because the target model does not support vision.]",
                            })
                    else:
                        logging.warning("Video frame extraction returned no frames. filename=%s", attachment.filename)
                        content_parts.append({"type": "text", "text": f"[No video frames extracted: {attachment.filename}]"})
                except Exception as e:
                    logging.exception("Video frame extraction failed. filename=%s", attachment.filename)
                    content_parts.append({"type": "text", "text": f"[Video frame extraction failed: {attachment.filename} - {str(e)[:80]}]"})

            if not content_parts:
                duration_text = self._format_seconds(duration) if duration is not None else "unknown"
                logging.info("Video processing produced metadata only. filename=%s duration=%s", attachment.filename, duration_text)
                return f"--- Video Metadata: {attachment.filename} ---\nDuration: {duration_text}\nFile Size: {file_size_mb:.2f} MB"

            logging.info(
                "Video attachment processing finished. filename=%s parts=%s summary=%s",
                attachment.filename,
                len(content_parts),
                self._processed_content_summary(content_parts),
            )
            return content_parts
        finally:
            self._remove_temp_file(temp_path)

    async def _process_attachment(self, attachment: discord.Attachment, model_name: str, media_settings: dict):
        mime_type = mimetypes.guess_type(attachment.filename)[0]
        file_ext = attachment.filename.lower().split('.')[-1] if '.' in attachment.filename else ''
        
        try:
            # Check file size limits first
            file_size_mb = getattr(attachment, "size", 0) / (1024 * 1024)
            attachment_url = getattr(attachment, "url", "")
            actual_content_type = ""
            logging.info(
                "Attachment processing started. filename=%s declared_mime=%s ext=%s declared_size_mb=%.2f model=%s",
                getattr(attachment, "filename", "unknown"),
                mime_type,
                file_ext,
                file_size_mb,
                model_name,
            )
            
            # Enhanced detection for GIFs, especially Tenor GIFs
            is_likely_gif = (
                mime_type == "image/gif" or 
                file_ext == "gif" or 
                '.gif' in attachment_url.lower() or
                'tenor.com' in attachment_url.lower() or
                'giphy.com' in attachment_url.lower() or
                '?format=gif' in attachment_url.lower()
            )
            logging.debug(
                "Attachment initial media detection. filename=%s is_likely_gif=%s supports_vision=%s",
                attachment.filename,
                is_likely_gif,
                self.llm_provider.supports_vision(model_name),
            )
            
            # Configure headers for better compatibility with GIF services
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': 'image/gif,image/webp,image/apng,image/*,*/*;q=0.8',
                'Accept-Encoding': 'gzip, deflate',
                'Connection': 'keep-alive',
            }

            if not attachment_url:
                return f"[Could not fetch attachment: {attachment.filename} (missing URL)]"

            safety_error = await self._validate_public_media_url(attachment_url)
            if safety_error:
                logging.warning(
                    "Rejected unsafe attachment URL. filename=%s url=%s reason=%s",
                    attachment.filename,
                    attachment_url,
                    safety_error,
                )
                return f"[Attachment URL rejected: {attachment.filename} - {safety_error}]"

            max_download_bytes = self._attachment_download_limit_bytes(media_settings, mime_type, file_ext, is_likely_gif)
            max_download_mb = max_download_bytes / (1024 * 1024)
            
            async with aiohttp.ClientSession() as session:
                logging.info("Downloading attachment. filename=%s", attachment.filename)
                async with session.get(attachment_url, headers=headers, allow_redirects=True, timeout=30) as resp:
                    final_url = str(resp.url)
                    safety_error = await self._validate_public_media_url(final_url)
                    if safety_error:
                        logging.warning(
                            "Rejected unsafe attachment redirect. filename=%s original_url=%s final_url=%s reason=%s",
                            attachment.filename,
                            attachment_url,
                            final_url,
                            safety_error,
                        )
                        return f"[Attachment redirect rejected: {attachment.filename} - {safety_error}]"

                    if resp.status != 200:
                        logging.warning("HTTP %s when fetching attachment. filename=%s", resp.status, attachment.filename)
                        return f"[Could not fetch attachment: {attachment.filename} (HTTP {resp.status})]"
                    
                    # Check actual content type from response headers
                    actual_content_type = resp.headers.get('content-type', '').lower()
                    actual_mime_type = actual_content_type.split(';', 1)[0].strip()
                    content_length = resp.headers.get('content-length')
                    response_size = None
                    if content_length:
                        try:
                            response_size = int(content_length)
                            if response_size > max_download_bytes:
                                logging.warning(
                                    "Attachment response over download limit. filename=%s content_length=%s max_download_mb=%.2f",
                                    attachment.filename,
                                    response_size,
                                    max_download_mb,
                                )
                                return f"[Attachment too large for processing: {attachment.filename} - larger than {max_download_mb:.1f}MB]"
                        except ValueError:
                            logging.debug("Attachment content-length was not numeric. filename=%s content_length=%s", attachment.filename, content_length)
                    
                    logging.debug(
                        "Fetched attachment headers. filename=%s content_type=%s content_length=%s",
                        attachment.filename,
                        actual_content_type,
                        content_length,
                    )
                    
                    file_bytes, exceeded_download_limit = await self._read_bounded_response_body(resp, max_download_bytes)
                    if exceeded_download_limit:
                        logging.warning(
                            "Attachment response exceeded bounded read. filename=%s downloaded_bytes=%s max_download_mb=%.2f",
                            attachment.filename,
                            len(file_bytes),
                            max_download_mb,
                        )
                        return f"[Attachment too large for processing: {attachment.filename} - larger than {max_download_mb:.1f}MB]"
                    if response_size is not None and not resp.headers.get("content-encoding") and len(file_bytes) < response_size:
                        logging.warning(
                            "Attachment download ended before content-length. filename=%s expected_bytes=%s downloaded_bytes=%s",
                            attachment.filename,
                            response_size,
                            len(file_bytes),
                        )
                        return f"[Could not fetch attachment: {attachment.filename} (incomplete download)]"
                    logging.info(
                        "Attachment downloaded. filename=%s status=%s actual_content_type=%s content_length_header=%s downloaded_bytes=%s",
                        attachment.filename,
                        resp.status,
                        actual_content_type,
                        content_length,
                        len(file_bytes),
                    )
                    
                    # Update file size based on actual downloaded content
                    if len(file_bytes) > 0:
                        file_size_mb = len(file_bytes) / (1024 * 1024)
                    
                    # Update MIME type based on actual response if it's more specific
                    if actual_mime_type and actual_mime_type != "application/octet-stream":
                        mime_type = actual_mime_type
                    
                    # Enhanced GIF detection with actual content type
                    if not is_likely_gif and actual_mime_type == 'image/gif':
                        is_likely_gif = True
                    logging.info(
                        "Attachment MIME finalized. filename=%s mime_type=%s actual_mime=%s size_mb=%.2f is_likely_gif=%s",
                        attachment.filename,
                        mime_type,
                        actual_mime_type,
                        file_size_mb,
                        is_likely_gif,
                    )

            # Image processing (including GIFs)
            if (mime_type and mime_type.startswith("image/")) or is_likely_gif:
                logging.info("Attachment routed to image processor. filename=%s mime_type=%s", attachment.filename, mime_type)
                images_config = media_settings.get("images", {})
                if not images_config.get("enabled", True):
                    logging.info("Image processing disabled. filename=%s", attachment.filename)
                    return f"[Image processing disabled: {attachment.filename}]"
                
                max_size = images_config.get("max_size_mb", 10)
                if file_size_mb > max_size:
                    logging.warning(
                        "Image attachment over size limit. filename=%s size_mb=%.2f max_size_mb=%s",
                        attachment.filename,
                        file_size_mb,
                        max_size,
                    )
                    return f"[Image too large for processing: {attachment.filename} - {file_size_mb:.1f}MB]"
                
                if self.llm_provider.supports_vision(model_name):
                    # Check if this is an animated GIF using enhanced detection
                    if is_likely_gif:
                        logging.info(
                            "Processing potential GIF. filename=%s content_type=%s mime=%s",
                            attachment.filename,
                            actual_content_type,
                            mime_type,
                        )
                        # Handle animated GIF by extracting frames
                        gif_config = images_config.get("gif", {})
                        if gif_config.get("extract_frames", True):
                            try:
                                import base64
                                
                                # Validate that we have image data before trying to process it
                                if len(file_bytes) == 0:
                                    logging.warning(f"No data received for {attachment.filename}")
                                    return self._image_content_part(attachment.url, model_name=model_name, filename=attachment.filename)
                                
                                gif = Image.open(io.BytesIO(file_bytes))
                                
                                # Check if it's actually animated
                                is_animated = getattr(gif, 'is_animated', False)
                                logging.info(f"GIF analysis for {attachment.filename}: is_animated={is_animated}, format={gif.format}")
                                if is_animated:
                                    # Extract frames from animated GIF with equal distribution across entire length
                                    max_frames = gif_config.get("max_frames", self.config.DEFAULT_GIF_MAX_FRAMES)
                                    frame_quality = gif_config.get("frame_quality", self.config.DEFAULT_GIF_FRAME_QUALITY)
                                    
                                    frames = []
                                    total_frames = getattr(gif, 'n_frames', 1)
                                    
                                    # Calculate frame indices for equal distribution across the entire GIF
                                    if total_frames <= max_frames:
                                        # If GIF has fewer frames than max, take all frames
                                        frame_indices = list(range(total_frames))
                                    else:
                                        # Distribute frames evenly across the entire GIF length
                                        # This ensures we cover the beginning, middle, and end of the animation
                                        if max_frames == 1:
                                            frame_indices = [0]
                                        elif max_frames == 2:
                                            frame_indices = [0, total_frames - 1]
                                        else:
                                            # For 3+ frames, distribute evenly including first and last frames
                                            step = (total_frames - 1) / (max_frames - 1)
                                            frame_indices = [round(i * step) for i in range(max_frames)]
                                            # Ensure the last frame is exactly the last frame
                                            frame_indices[-1] = total_frames - 1
                                            # Remove duplicates while preserving order
                                            seen = set()
                                            frame_indices = [x for x in frame_indices if not (x in seen or seen.add(x))]
                                    
                                    logging.info(f"Processing animated GIF {attachment.filename}: {total_frames} total frames, extracting frames at indices {frame_indices}")
                                    
                                    try:
                                        ocr_frame_bytes = []
                                        for frame_index in frame_indices:
                                            try:
                                                gif.seek(frame_index)
                                                # Convert frame to RGB (remove transparency)
                                                frame = gif.convert('RGB')
                                                
                                                # Convert frame to base64 for the API
                                                buffer = io.BytesIO()
                                                frame.save(buffer, format='JPEG', quality=frame_quality)
                                                buffer.seek(0)
                                                frame_bytes = buffer.getvalue()
                                                frame_b64 = base64.b64encode(frame_bytes).decode()
                                                frames.append(f"data:image/jpeg;base64,{frame_b64}")
                                                ocr_frame_bytes.append({
                                                    "label": f"frame {frame_index}",
                                                    "jpeg_bytes": frame_bytes,
                                                })
                                                
                                            except (EOFError, OSError) as e:
                                                logging.warning(f"Could not access frame {frame_index} in {attachment.filename}: {e}")
                                                continue
                                                
                                    except Exception as frame_error:
                                        logging.warning(f"Error extracting frames from {attachment.filename}: {frame_error}")
                                        # Try fallback method with sequential frame access
                                        frames = []
                                        ocr_frame_bytes = []
                                        try:
                                            frame_count = 0
                                            current_frame = 0
                                            while frame_count < max_frames and current_frame < total_frames:
                                                # Convert frame to RGB (remove transparency)
                                                frame = gif.convert('RGB')
                                                
                                                # Convert frame to base64 for the API
                                                buffer = io.BytesIO()
                                                frame.save(buffer, format='JPEG', quality=frame_quality)
                                                buffer.seek(0)
                                                frame_bytes = buffer.getvalue()
                                                frame_b64 = base64.b64encode(frame_bytes).decode()
                                                frames.append(f"data:image/jpeg;base64,{frame_b64}")
                                                ocr_frame_bytes.append({
                                                    "label": f"frame {current_frame}",
                                                    "jpeg_bytes": frame_bytes,
                                                })
                                                
                                                frame_count += 1
                                                current_frame += max(1, total_frames // max_frames)
                                                gif.seek(current_frame)
                                        except (EOFError, OSError):
                                            # End of GIF frames
                                            pass
                                    
                                    if frames:
                                        logging.info(
                                            "Animated GIF frames extracted. filename=%s total_frames=%s extracted_frames=%s",
                                            attachment.filename,
                                            total_frames,
                                            len(frames),
                                        )
                                        frame_ocr_text = await asyncio.to_thread(
                                            self._ocr_image_frames_text_for_context,
                                            ocr_frame_bytes,
                                            attachment.filename,
                                            images_config,
                                        )
                                        # Return structured data for animated GIFs that can be used in different contexts
                                        return {
                                            "type": "animated_gif",
                                            "total_frames": total_frames,
                                            "extracted_frames": len(frames),
                                            "frames": frames,
                                            "filename": attachment.filename,
                                            "ocr_text": frame_ocr_text,
                                        }
                                    else:
                                        # Fallback to treating as static image
                                        logging.warning("Animated GIF extraction produced no frames; falling back to image. filename=%s", attachment.filename)
                                        return await asyncio.to_thread(
                            self._vision_image_content_parts,
                                            url=attachment.url,
                                            file_bytes=file_bytes,
                                            mime_type=mime_type,
                                            model_name=model_name,
                                            filename=attachment.filename,
                                            images_config=images_config,
                                        )
                                else:
                                    # Static GIF, treat as regular image
                                    logging.info("GIF is not animated; sending as static image. filename=%s", attachment.filename)
                                    return await asyncio.to_thread(
                            self._vision_image_content_parts,
                                        url=attachment.url,
                                        file_bytes=file_bytes,
                                        mime_type=mime_type,
                                        model_name=model_name,
                                        filename=attachment.filename,
                                        images_config=images_config,
                                    )
                            except Exception as e:
                                logging.warning(
                                    "Failed to process potential GIF. filename=%s error_type=%s error=%s",
                                    attachment.filename,
                                    type(e).__name__,
                                    e,
                                )
                                # Check if this was an image format issue
                                if "cannot identify image file" in str(e).lower() or "truncated" in str(e).lower():
                                    logging.info(f"Image format issue with {attachment.filename}, might not be a valid image file")
                                # Fallback to treating as static image
                                return self._image_content_part(attachment.url, model_name=model_name, filename=attachment.filename)
                        else:
                            # GIF frame extraction disabled, treat as static image
                            logging.info("GIF frame extraction disabled; sending as static image. filename=%s", attachment.filename)
                            return await asyncio.to_thread(
                            self._vision_image_content_parts,
                                url=attachment.url,
                                file_bytes=file_bytes,
                                mime_type=mime_type,
                                model_name=model_name,
                                filename=attachment.filename,
                                images_config=images_config,
                            )
                    else:
                        # Regular static image
                        logging.info("Static image will be attached to vision model. filename=%s model=%s", attachment.filename, model_name)
                        return await asyncio.to_thread(
                            self._vision_image_content_parts,
                            url=attachment.url,
                            file_bytes=file_bytes,
                            mime_type=mime_type,
                            model_name=model_name,
                            filename=attachment.filename,
                            images_config=images_config,
                        )
                else:
                    # Fallback: OCR for text-only models
                    logging.info("Image target model has no vision; attempting OCR fallback. filename=%s model=%s", attachment.filename, model_name)
                    if images_config.get("ocr_enabled", True) and TESSERACT_AVAILABLE:
                        try:
                            with Image.open(io.BytesIO(file_bytes)) as img:
                                ocr_text = await asyncio.to_thread(pytesseract.image_to_string, img)
                                if ocr_text.strip():
                                    logging.info("Image OCR succeeded. filename=%s chars=%s", attachment.filename, len(ocr_text.strip()))
                                    return f"--- Image OCR Text: {attachment.filename} ---\n{ocr_text.strip()}"
                                else:
                                    logging.info("Image OCR found no readable text. filename=%s", attachment.filename)
                                    return f"[Image contains no readable text: {attachment.filename}]"
                        except Exception as e:
                            logging.exception("Image OCR failed. filename=%s", attachment.filename)
                            if images_config.get("description_fallback", True):
                                return f"[Image with OCR failure: {attachment.filename} - {str(e)[:50]}...]"
                            else:
                                return f"[Image omitted: {attachment.filename}]"
                    else:
                        logging.info(
                            "Image omitted for text-only model because OCR is disabled or unavailable. filename=%s ocr_enabled=%s tesseract_available=%s",
                            attachment.filename,
                            images_config.get("ocr_enabled", True),
                            TESSERACT_AVAILABLE,
                        )
                        return f"[Image omitted - OCR disabled or unavailable: {attachment.filename}]"

            # Audio processing
            elif self._is_audio_file(mime_type, file_ext):
                logging.info("Attachment routed to audio processor. filename=%s mime_type=%s", attachment.filename, mime_type)
                audio_config = media_settings.get("audio", {})
                return await self._process_audio_attachment(
                    attachment,
                    file_bytes,
                    file_ext,
                    file_size_mb,
                    model_name,
                    audio_config,
                )

            # Video processing
            elif self._is_video_file(mime_type, file_ext):
                logging.info("Attachment routed to video processor. filename=%s mime_type=%s", attachment.filename, mime_type)
                video_config = media_settings.get("video", {})
                return await self._process_video_attachment(
                    attachment,
                    file_bytes,
                    file_ext,
                    file_size_mb,
                    model_name,
                    video_config,
                )

            # PDF processing
            elif mime_type == "application/pdf":
                logging.info("Attachment routed to PDF processor. filename=%s", attachment.filename)
                pdf_config = media_settings.get("pdf", {})
                if not pdf_config.get("enabled", True):
                    logging.info("PDF processing disabled. filename=%s", attachment.filename)
                    return f"[PDF processing disabled: {attachment.filename}]"
                
                max_size = pdf_config.get("max_size_mb", 10)
                if file_size_mb > max_size:
                    logging.warning(
                        "PDF attachment over size limit. filename=%s size_mb=%.2f max_size_mb=%s",
                        attachment.filename,
                        file_size_mb,
                        max_size,
                    )
                    return f"[PDF too large for processing: {attachment.filename} - {file_size_mb:.1f}MB]"
                
                if self.llm_provider.supports_pdf(model_name):
                    # Send directly as PDF data for PDF-capable models
                    logging.info("Sending PDF directly to model. filename=%s model=%s", attachment.filename, model_name)
                    return {"type": "pdf_url", "pdf_url": {"url": attachment.url}}
                else:
                    # Fallback: Text extraction for all models
                    try:
                        text = ""
                        if PDFPLUMBER_AVAILABLE:
                            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                                for i, page in enumerate(pdf.pages):
                                    page_text = page.extract_text() or ""
                                    if pdf_config.get("preserve_formatting", True):
                                        text += f"\n--- Page {i + 1} ---\n{page_text}\n"
                                    else:
                                        text += page_text + "\n"
                        else: # Fallback to PyPDF2
                            with io.BytesIO(file_bytes) as f:
                                reader = PyPDF2.PdfReader(f)
                                for page_num, page in enumerate(reader.pages):
                                    page_text = page.extract_text() or ""
                                    if pdf_config.get("preserve_formatting", True):
                                        text += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
                                    else:
                                        text += page_text + "\n"
                        logging.info("PDF text extracted. filename=%s chars=%s", attachment.filename, len(text.strip()))
                        return f"--- PDF Content: {attachment.filename} ---\n{self._truncate_document_text(text, pdf_config)}"
                    except Exception as e:
                        logging.exception("PDF text extraction failed. filename=%s", attachment.filename)
                        return f"[PDF text extraction failed: {attachment.filename} - {str(e)[:50]}...]"

            # Office Documents (DOCX, XLSX, PPTX)
            elif mime_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                               "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", 
                               "application/vnd.openxmlformats-officedocument.presentationml.presentation"] or file_ext in OFFICE_EXTENSIONS:
                logging.info("Attachment routed to Office document processor. filename=%s ext=%s", attachment.filename, file_ext)
                office_config = media_settings.get("office_documents", {})
                if not office_config.get("enabled", True):
                    logging.info("Office document processing disabled. filename=%s", attachment.filename)
                    return f"[Office document processing disabled: {attachment.filename}]"
                
                max_size = office_config.get("max_size_mb", 10)
                if file_size_mb > max_size:
                    logging.warning(
                        "Office document over size limit. filename=%s size_mb=%.2f max_size_mb=%s",
                        attachment.filename,
                        file_size_mb,
                        max_size,
                    )
                    return f"[Document too large for processing: {attachment.filename} - {file_size_mb:.1f}MB]"
                
                preserve_structure = office_config.get("preserve_structure", True)
                
                try:
                    if file_ext == "docx":
                        doc = docx.Document(io.BytesIO(file_bytes))
                        if preserve_structure:
                            text = ""
                            for para in doc.paragraphs:
                                if para.style.name.startswith('Heading'):
                                    text += f"\n## {para.text}\n"
                                else:
                                    text += f"{para.text}\n"
                        else:
                            text = "\n".join([para.text for para in doc.paragraphs])
                        logging.info("DOCX text extracted. filename=%s paragraphs=%s chars=%s", attachment.filename, len(doc.paragraphs), len(text.strip()))
                        return f"--- Document Content: {attachment.filename} ---\n{self._truncate_document_text(text, office_config)}"
                    
                    elif file_ext == "xlsx":
                        wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
                        text = ""
                        for sheet_name in wb.sheetnames:
                            sheet = wb[sheet_name]
                            if preserve_structure:
                                text += f"\n--- Sheet: {sheet_name} ---\n"
                            for row in sheet.iter_rows(values_only=True):
                                text += ", ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
                        logging.info("XLSX text extracted. filename=%s sheets=%s chars=%s", attachment.filename, len(wb.sheetnames), len(text.strip()))
                        return f"--- Excel Content: {attachment.filename} ---\n{self._truncate_document_text(text, office_config)}"
                    
                    elif file_ext == "pptx" and PPTX_AVAILABLE:
                        prs = pptx.Presentation(io.BytesIO(file_bytes))
                        text = ""
                        for i, slide in enumerate(prs.slides):
                            if preserve_structure:
                                text += f"\n--- Slide {i+1} ---\n"
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text += shape.text + "\n"
                        logging.info("PPTX text extracted. filename=%s slides=%s chars=%s", attachment.filename, len(prs.slides), len(text.strip()))
                        return f"--- PowerPoint Content: {attachment.filename} ---\n{self._truncate_document_text(text, office_config)}"
                    elif file_ext == "pptx" and not PPTX_AVAILABLE:
                        logging.info("PPTX processing unavailable because python-pptx is not installed. filename=%s", attachment.filename)
                        return f"[PowerPoint processing unavailable - python-pptx not installed: {attachment.filename}]"
                        
                except Exception as e:
                    logging.exception("Office document processing failed. filename=%s", attachment.filename)
                    return f"[Document processing failed: {attachment.filename} - {str(e)[:50]}...]"

            # Text files
            elif (mime_type and mime_type.startswith("text/")) or file_ext in TEXT_EXTENSIONS:
                logging.info("Attachment routed to text file processor. filename=%s ext=%s", attachment.filename, file_ext)
                text_config = media_settings.get("text_files", {})
                if not text_config.get("enabled", True):
                    logging.info("Text file processing disabled. filename=%s", attachment.filename)
                    return f"[Text file processing disabled: {attachment.filename}]"
                
                max_size = text_config.get("max_size_mb", 5)
                if file_size_mb > max_size:
                    logging.warning(
                        "Text file over size limit. filename=%s size_mb=%.2f max_size_mb=%s",
                        attachment.filename,
                        file_size_mb,
                        max_size,
                    )
                    return f"[Text file too large for processing: {attachment.filename} - {file_size_mb:.1f}MB]"
                
                supported_extensions = text_config.get("supported_extensions", sorted(TEXT_EXTENSIONS))
                if file_ext not in supported_extensions:
                    logging.info("Text file extension unsupported. filename=%s ext=%s", attachment.filename, file_ext)
                    return f"[Text file extension not supported: {attachment.filename}]"
                
                try:
                    # Full content read and included as text context for both decision and main models
                    text_content = file_bytes.decode('utf-8', errors='ignore')
                    logging.info("Text file decoded. filename=%s chars=%s", attachment.filename, len(text_content))
                    return f"--- File Content: {attachment.filename} ---\n{self._truncate_document_text(text_content, text_config)}"
                except Exception as e:
                    logging.exception("Text file reading failed. filename=%s", attachment.filename)
                    return f"[Text file reading failed: {attachment.filename} - {str(e)[:50]}...]"

            # Other files - metadata only
            else:
                logging.info("Attachment routed to fallback/other processor. filename=%s mime_type=%s ext=%s", attachment.filename, mime_type, file_ext)
                other_config = media_settings.get("other_files", {})
                if not other_config.get("enabled", True):
                    logging.info("Other file processing disabled. filename=%s", attachment.filename)
                    return f"[Other file processing disabled: {attachment.filename}]"
                
                max_size = other_config.get("max_size_mb", 20)
                if file_size_mb > max_size:
                    logging.warning(
                        "Other file over size limit. filename=%s size_mb=%.2f max_size_mb=%s",
                        attachment.filename,
                        file_size_mb,
                        max_size,
                    )
                    return f"[File too large for processing: {attachment.filename} - {file_size_mb:.1f}MB]"
                
                if other_config.get("include_metadata_only", True):
                    # File name, size, type, and basic metadata included as text description
                    logging.info("Returning metadata for unsupported attachment. filename=%s", attachment.filename)
                    return f"--- File Metadata: {attachment.filename} ---\nFile Type: {mime_type or 'Unknown'}\nFile Size: {file_size_mb:.2f} MB\nFile Extension: .{file_ext}\nUpload URL: {attachment.url}"
                else:
                    logging.info("Unsupported attachment omitted. filename=%s", attachment.filename)
                    return f"[Unsupported file type: {attachment.filename}]"

        except Exception as e:
            logging.exception(f"Error processing attachment {attachment.filename}: {e}")
            return f"[Could not process file: {attachment.filename} - Error: {str(e)[:50]}...]"

    async def _process_embed(self, embed: discord.Embed, model_name: str, media_settings: dict):
        """
        Process Discord embeds, with special handling for Tenor GIFs and other media.
        """
        try:
            # Check if this embed contains media we can process
            embed_url = embed.url
            embed_image_url = embed.image.url if embed.image else None
            embed_video_url = embed.video.url if embed.video else None
            
            # Enhanced detection for embedded GIFs (especially Tenor)
            potential_gif_urls = []
            if embed_image_url:
                potential_gif_urls.append(embed_image_url)
            if embed_video_url:
                potential_gif_urls.append(embed_video_url)
            if embed_url:
                potential_gif_urls.append(embed_url)
            
            for url in potential_gif_urls:
                if url and (
                    'tenor.com' in url.lower() or 
                    'giphy.com' in url.lower() or 
                    '.gif' in url.lower() or
                    url.endswith('.gif')
                ):
                    logging.info(f"Processing embedded GIF from URL: {url}")
                    
                    # Create a fake attachment-like object for processing
                    class EmbedAttachment:
                        def __init__(self, url, filename=None):
                            self.url = url
                            self.filename = filename or url.split('/')[-1] or 'embedded_gif.gif'
                            self.size = 0  # Unknown size for embeds
                    
                    fake_attachment = EmbedAttachment(url, f"embedded_{embed.type or 'media'}.gif")
                    
                    # Process using the same logic as attachments
                    try:
                        processed_content = await self._process_attachment(fake_attachment, model_name, media_settings)
                        if processed_content:
                            return processed_content
                    except Exception as e:
                        logging.warning(f"Failed to process embedded media from {url}: {e}")
                        continue
            
            # If no media was processed, return basic embed information
            embed_info = []
            if embed.title:
                embed_info.append(f"Embed Title: {embed.title}")
            if embed.description:
                embed_info.append(f"Embed Description: {embed.description[:200]}...")
            if embed.url and not any(url in potential_gif_urls for url in [embed.url]):
                embed_info.append(f"Embed URL: {embed.url}")
            
            if embed_info:
                return "--- Embedded Content ---\n" + "\n".join(embed_info)
            
            return None  # No processable content found
            
        except Exception as e:
            logging.warning(f"Error processing embed: {e}")
            return f"[Could not process embedded content - Error: {str(e)[:50]}...]"

    async def update_channel_summary(self, guild_id: str, channel_id: str):
        """
        Update the channel summary by integrating recent messages with the existing summary.
        This uses an incremental approach to build a comprehensive, evolving narrative.
        """
        logging.info(f"Updating summary for channel {channel_id} in guild {guild_id}...")
        
        try:
            data = await self.store.get_guild_data(guild_id)
            settings = await self.get_guild_and_channel_settings(guild_id, channel_id)
            
            # Get channel and validate it exists
            channel = self.bot.get_channel(int(channel_id))
            if not channel:
                logging.warning(f"Cannot update summary, channel {channel_id} not found.")
                return
            
            # Get current channel data
            channel_data = data.get("channels", {}).get(channel_id, {})
            old_summary = channel_data.get("summary", "")
            last_summary_time = channel_data.get("last_summary_time")
            
            # Determine if this is the first summary
            is_first_summary = not old_summary or old_summary.strip() == "" or old_summary == "No summary yet."
            
            # Determine how many messages to fetch based on settings
            if is_first_summary:
                history_limit = settings.get("initial_summarize_messages", self.config.INITIAL_SUMMARY_MESSAGES)
            else:
                history_limit = settings.get("summarize_every_messages", self.config.DEFAULT_SUMMARIZE_EVERY_MESSAGES)
            
            # If we have a last summary time, try to fetch messages since then
            # Otherwise, use the message count limit
            after_time = None
            if last_summary_time:
                try:
                    after_time = datetime.fromisoformat(last_summary_time)
                    # Add a small buffer to avoid missing messages due to timing
                    after_time = after_time.replace(microsecond=0)
                except ValueError:
                    logging.warning(f"Invalid last_summary_time format: {last_summary_time}")
            
            # Fetch recent messages
            recent_messages = []
            message_count = 0
            
            async for msg in channel.history(limit=min(history_limit * 2, 500), after=after_time):
                # Filter out bot messages (except for important context)
                if msg.author.bot and msg.author.id != self.bot.user.id:
                    continue
                
                # Skip very short messages that don't add context
                if len(msg.content.strip()) < 3 and not msg.attachments:
                    continue
                
                recent_messages.append(msg)
                message_count += 1
                
                # If we have enough messages and no specific time constraint, stop
                if not after_time and message_count >= history_limit:
                    break
            
            # channel.history yields newest-first by default but oldest-first when
            # `after` is set, so only the no-after fetch needs reversing.
            if not after_time:
                recent_messages.reverse()  # Oldest to newest

            if not recent_messages:
                logging.info(f"No new messages to summarize for channel {channel_id}.")
                return
            
            logging.info(f"Processing {len(recent_messages)} messages for channel summary update.")
            
            # Build conversation text with rich context
            conversation_lines = []
            for msg in recent_messages:
                # Format message with timestamp and user info
                timestamp = msg.created_at.strftime("%Y-%m-%d %H:%M")
                author_name = msg.author.display_name
                author_id = msg.author.id
                
                # Handle message content
                content = msg.content.strip() if msg.content else ""
                
                # Add attachment information
                attachment_info = ""
                if msg.attachments:
                    attachments = [f"[{att.filename}]" for att in msg.attachments]
                    attachment_info = f" (attachments: {', '.join(attachments)})"
                
                # Handle replies
                reply_info = ""
                if msg.reference and msg.reference.resolved:
                    if hasattr(msg.reference.resolved, 'author'):
                        try:
                            reply_to = msg.reference.resolved.author.display_name
                            reply_to_id = msg.reference.resolved.author.id
                            reply_info = f" (replying to {reply_to} [{reply_to_id}])"
                        except AttributeError:
                            # Handle deleted accounts
                            reply_info = " (replying to [Deleted Account])"
                    else:
                        # Handle deleted messages
                        reply_info = " (replying to [Deleted Message])"
                
                # Combine all parts
                if content or attachment_info:
                    full_message = f"[{timestamp}] {author_name} [{author_id}]{reply_info}: {content}{attachment_info}"
                    conversation_lines.append(full_message)
            
            conversation_text = "\n".join(conversation_lines)
            
            # Determine if this is the first summary or an update
            is_first_summary = not old_summary or old_summary.strip() == "" or old_summary == "No summary yet."
            
            if is_first_summary:
                # Create initial summary
                prompt = f"""You are a conversation summarizer for a Discord channel. 
Analyze the following conversation and create a comprehensive summary that captures:
- Main topics and themes discussed
- Key questions asked and answers provided
- Important decisions or conclusions reached
- Notable events or announcements
- Active participants and their contributions

The summary should be well-organized, informative, and provide a clear overview of what this channel is about and what has been discussed.

Channel: #{channel.name}
Server: {channel.guild.name}

Conversation History:
---
{conversation_text}
---

Provide a comprehensive channel summary:"""
            else:
                # Update existing summary incrementally
                prompt = f"""You are a conversation summarizer for a Discord channel.
You have an existing summary of the channel's conversation history. Now you need to integrate new messages into this summary.

Current Channel Summary:
---
{old_summary}
---

New Messages to Integrate:
---
{conversation_text}
---

Instructions:
1. Review the existing summary and the new messages
2. Identify new topics, developments, or conclusions from the recent messages
3. Update the existing summary by integrating the new information
4. Maintain the chronological flow and thematic organization
5. Remove or update outdated information if necessary
6. Keep the summary comprehensive but concise

Provide the updated, complete channel summary:"""
            
            # Generate the summary using the LLM
            response = await self.llm_provider.create_completion(
                model=settings.get("model", self.config.MAIN_LLM_MODEL),
                messages=[
                    {"role": "system", "content": "You are an expert conversation summarizer. Create clear, comprehensive, and well-organized summaries that capture the essence of Discord channel conversations."},
                    {"role": "user", "content": prompt}
                ]
            )
            
            if not response or not response.choices:
                logging.error(f"Failed to get a valid response from LLM for channel summary {channel_id}.")
                return
            
            new_summary = response.choices[0].message.content.strip()
            
            # Validate the summary isn't empty or generic
            if len(new_summary) < 50 or new_summary.lower().startswith("i cannot") or new_summary.lower().startswith("i can't"):
                logging.warning(f"Generated summary for channel {channel_id} seems invalid or too short: {new_summary[:100]}...")
                return
            
            # Update data file with new summary and metadata
            fresh_data = await self.store.get_data()
            if str(guild_id) not in fresh_data:
                fresh_data[str(guild_id)] = {}
            if "channels" not in fresh_data[str(guild_id)]:
                fresh_data[str(guild_id)]["channels"] = {}
            if channel_id not in fresh_data[str(guild_id)]["channels"]:
                fresh_data[str(guild_id)]["channels"][channel_id] = {}
            
            # Store the new summary and reset counters
            fresh_data[str(guild_id)]["channels"][channel_id]["summary"] = new_summary
            fresh_data[str(guild_id)]["channels"][channel_id]["messages_since_summary"] = 0
            fresh_data[str(guild_id)]["channels"][channel_id]["last_summary_time"] = datetime.now(timezone.utc).isoformat()
            fresh_data[str(guild_id)]["channels"][channel_id]["messages_processed"] = len(recent_messages)
            fresh_data[str(guild_id)]["channels"][channel_id]["summary_type"] = "initial" if is_first_summary else "incremental"
            
            # Save the updated data
            await self.store.save_data(fresh_data)
            
            summary_type = "Initial" if is_first_summary else "Incremental"
            logging.info(f"Successfully updated {summary_type.lower()} summary for channel #{channel.name} ({channel_id}). Processed {len(recent_messages)} messages.")
            
        except Exception as e:
            logging.error(f"Error updating channel summary for {channel_id}: {e}")
            # Don't re-raise to avoid breaking the message processing flow

    async def get_channel_summary(self, guild_id: str, channel_id: str) -> dict:
        """
        Get the current channel summary and metadata.
        
        Returns:
            dict: Contains summary text, last update time, message count, etc.
        """
        try:
            data = await self.store.get_guild_data(guild_id)
            channel_data = data.get("channels", {}).get(channel_id, {})
            
            return {
                "summary": channel_data.get("summary", "No summary available yet."),
                "last_summary_time": channel_data.get("last_summary_time"),
                "messages_since_summary": channel_data.get("messages_since_summary", 0),
                "messages_processed": channel_data.get("messages_processed", 0),
                "summary_type": channel_data.get("summary_type", "none")
            }
        except Exception as e:
            logging.error(f"Error getting channel summary for {channel_id}: {e}")
            return {
                "summary": "Error retrieving summary.",
                "last_summary_time": None,
                "messages_since_summary": 0,
                "messages_processed": 0,
                "summary_type": "error"
            }

    async def force_channel_summary_update(self, guild_id: str, channel_id: str) -> bool:
        """
        Force an immediate update of the channel summary regardless of message count or time.
        
        Returns:
            bool: True if successful, False otherwise
        """
        try:
            logging.info(f"Forcing channel summary update for {channel_id} in guild {guild_id}")
            await self.update_channel_summary(guild_id, channel_id)
            return True
        except Exception as e:
            logging.error(f"Error forcing channel summary update for {channel_id}: {e}")
            return False

    async def clear_channel_summary(self, guild_id: str, channel_id: str) -> bool:
        """
        Clear the channel summary and reset counters.
        
        Returns:
            bool: True if successful, False otherwise
        """
        try:
            fresh_data = await self.store.get_data()
            
            if str(guild_id) not in fresh_data:
                fresh_data[str(guild_id)] = {}
            if "channels" not in fresh_data[str(guild_id)]:
                fresh_data[str(guild_id)]["channels"] = {}
            if channel_id not in fresh_data[str(guild_id)]["channels"]:
                return False  # No summary to clear
            
            # Clear summary but keep the channel entry
            channel_data = fresh_data[str(guild_id)]["channels"][channel_id]
            channel_data.pop("summary", None)
            channel_data.pop("last_summary_time", None)
            channel_data["messages_since_summary"] = 0
            channel_data.pop("messages_processed", None)
            channel_data.pop("summary_type", None)
            
            await self.store.save_data(fresh_data)
            logging.info(f"Cleared channel summary for {channel_id}")
            return True
            
        except Exception as e:
            logging.error(f"Error clearing channel summary for {channel_id}: {e}")
            return False

    async def update_user_profile(self, guild_id: str, user_id: str, guild_obj=None):
        """
        Updates the AI-generated profile for a user by gathering their recent messages
        across all channels in the guild and generating/updating their summary.
        """
        logging.info(f"Updating AI profile for user {user_id} in guild {guild_id}...")
        
        try:
            # Get current data and settings
            data = await self.store.get_guild_data(guild_id)
            settings = await self.get_guild_and_channel_settings(guild_id, None)
            
            # Initialize guild data structure if needed
            if "users" not in data:
                data["users"] = {}
            if user_id not in data["users"]:
                data["users"][user_id] = {}
                
            user_data = data["users"][user_id]
            old_summary = user_data.get("ai_summary", "")
            manual_note = user_data.get("manual_note", "")
            
            # Gather recent messages from the user across all accessible channels
            user_messages = await self._gather_user_messages(guild_id, user_id, guild_obj)
            
            if not user_messages:
                logging.warning(f"No messages found for user {user_id} in guild {guild_id}")
                # Reset counters even if no messages found
                user_data["messages_since_profile_update"] = 0
                user_data["last_profile_update_time"] = datetime.now(timezone.utc).isoformat()
                
                # Get fresh data and update to avoid overwriting concurrent changes
                fresh_data = await self.store.get_data()
                if str(guild_id) not in fresh_data:
                    fresh_data[str(guild_id)] = {}
                if "users" not in fresh_data[str(guild_id)]:
                    fresh_data[str(guild_id)]["users"] = {}
                if user_id not in fresh_data[str(guild_id)]["users"]:
                    fresh_data[str(guild_id)]["users"][user_id] = {}
                
                fresh_data[str(guild_id)]["users"][user_id].update(user_data)
                await self.store.save_data(fresh_data)
                return
            
            # Generate AI summary using LLM
            new_summary = await self._generate_user_profile_summary(
                user_messages, old_summary, manual_note, user_id, settings
            )
            
            if new_summary:
                # Update user data
                user_data["ai_summary"] = new_summary
                user_data["messages_since_profile_update"] = 0
                user_data["last_profile_update_time"] = datetime.now(timezone.utc).isoformat()
                
                # Get fresh data and update to avoid overwriting concurrent changes
                fresh_data = await self.store.get_data()
                if str(guild_id) not in fresh_data:
                    fresh_data[str(guild_id)] = {}
                if "users" not in fresh_data[str(guild_id)]:
                    fresh_data[str(guild_id)]["users"] = {}
                if user_id not in fresh_data[str(guild_id)]["users"]:
                    fresh_data[str(guild_id)]["users"][user_id] = {}
                
                fresh_data[str(guild_id)]["users"][user_id].update(user_data)
                await self.store.save_data(fresh_data)
                logging.info(f"AI profile updated successfully for user {user_id}")
            else:
                logging.error(f"Failed to generate AI summary for user {user_id}")
                
        except Exception as e:
            logging.error(f"Error updating user profile for {user_id}: {e}")
            
    async def _gather_user_messages(self, guild_id: str, user_id: str, guild_obj=None, max_messages: int = 100):
        """
        Gathers recent messages from a user across all accessible channels in the guild.
        """
        if not guild_obj:
            # Try to get guild from bot if not provided
            try:
                guild_obj = discord.utils.get(self.bot.guilds, id=int(guild_id))
                if not guild_obj:
                    logging.error(f"Guild {guild_id} not found")
                    return []
            except Exception as e:
                logging.error(f"Error getting guild {guild_id}: {e}")
                return []
        
        user_messages = []
        channels_checked = 0
        
        try:
            # Get user object
            user = guild_obj.get_member(int(user_id))
            if not user:
                logging.warning(f"User {user_id} not found in guild {guild_id}")
                return []
            
            # Iterate through all text channels the bot can access
            for channel in guild_obj.text_channels:
                try:
                    # Check if bot has permission to read message history
                    if not channel.permissions_for(guild_obj.me).read_message_history:
                        continue
                        
                    channels_checked += 1
                    
                    # Fetch recent messages from this channel (limit per channel to avoid overload)
                    async for message in channel.history(limit=200):  # Look through more messages to find user's
                        if message.author.id == int(user_id):
                            if len(user_messages) >= max_messages:
                                break
                                
                            # Format message with context
                            formatted_message = {
                                'content': message.content or '[No text content]',
                                'channel': channel.name,
                                'timestamp': message.created_at.isoformat(),
                                'has_attachments': len(message.attachments) > 0,
                                'attachment_types': [att.content_type or 'unknown' for att in message.attachments] if message.attachments else [],
                                'embeds_count': len(message.embeds),
                                'reactions_count': len(message.reactions)
                            }
                            user_messages.append(formatted_message)
                            
                        if len(user_messages) >= max_messages:
                            break
                            
                except discord.Forbidden:
                    # Bot doesn't have permission to read this channel
                    continue
                except Exception as e:
                    logging.warning(f"Error reading messages from channel {channel.name}: {e}")
                    continue
                    
                if len(user_messages) >= max_messages:
                    break
            
            # Sort messages by timestamp (oldest first for chronological context)
            user_messages.sort(key=lambda x: x['timestamp'])
            
            logging.info(f"Gathered {len(user_messages)} messages from {channels_checked} channels for user {user_id}")
            return user_messages
            
        except Exception as e:
            logging.error(f"Error gathering messages for user {user_id}: {e}")
            return []
    
    async def _generate_user_profile_summary(self, user_messages, old_summary, manual_note, user_id, settings):
        """
        Generates or updates an AI summary of the user based on their messages.
        """
        try:
            model = settings.get("model", self.config.MAIN_LLM_MODEL)
            
            # Prepare message context (limit to prevent overly long prompts)
            message_context = []
            for msg in user_messages[-30:]:  # Reduced from 50 to 30 messages to keep prompt shorter
                context_line = f"[{msg['channel']}] {msg['content'][:200]}"  # Truncate long messages
                if msg['has_attachments']:
                    context_line += f" [Attachments: {len(msg['attachment_types'])} files]"
                if msg['embeds_count'] > 0:
                    context_line += f" [Embeds: {msg['embeds_count']}]"
                message_context.append(context_line)
            
            messages_text = "\n".join(message_context)
            
            # Build the prompt for profile generation/update (more concise)
            if old_summary and old_summary.strip():
                # Incremental update
                prompt = f"""Update this user's Discord profile summary with their recent messages.

CURRENT SUMMARY:
{old_summary}

RECENT MESSAGES:
{messages_text}

ADMIN NOTE: {manual_note or 'None'}

Provide an updated 2-paragraph summary focusing on communication style, interests, and notable traits."""
            else:
                # Initial profile generation
                prompt = f"""Create a Discord user profile summary from these messages.

MESSAGES:
{messages_text}

ADMIN NOTE: {manual_note or 'None'}

Provide a 2-paragraph summary covering communication style, interests, activity patterns, and key traits."""

            # Make the API call
            response = await self.llm_provider.create_completion(
                model=model,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=2000,  # Further increased token limit
                temperature=0.3
            )
            
            if response and response.choices and len(response.choices) > 0:
                choice = response.choices[0]
                if choice.message and choice.message.content:
                    return choice.message.content.strip()
                elif choice.finish_reason == 'length':
                    logging.error("LLM response was truncated due to token limit. Consider increasing max_tokens.")
                    return None
                else:
                    logging.error(f"LLM response has no content. Finish reason: {choice.finish_reason}")
                    return None
            else:
                logging.error("Empty or invalid response from LLM for user profile generation")
                return None
                
        except Exception as e:
            logging.error(f"Error generating user profile summary: {e}")
            return None
